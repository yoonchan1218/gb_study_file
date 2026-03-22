import copy
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.oxml.ns import qn
from lxml import etree

prs = Presentation("김윤찬_TRY-CATCH_화면_포트폴리오_v6.pptx")
NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
}

# ============================================================
# FIX 1: Slide 3 - Make green overlays semi-transparent
# ============================================================
slide3 = prs.slides[2]
shapes3 = list(slide3.shapes)

green_overlay_indices = []
for i, shape in enumerate(shapes3):
    sp = shape._element
    spPr = sp.find('.//p:spPr', NS)
    if spPr is None:
        spPr = sp.find('p:spPr', NS)
    if spPr is None:
        continue
    fill = spPr.find('a:solidFill', NS)
    if fill is None:
        continue
    clr = fill.find('a:srgbClr', NS)
    if clr is not None and clr.get('val') == 'DBEEE0':
        green_overlay_indices.append(i)
        # Add alpha for transparency
        # Remove existing alpha if any
        for existing_alpha in clr.findall('a:alpha', NS):
            clr.remove(existing_alpha)
        alpha = etree.SubElement(clr, qn('a:alpha'))
        alpha.set('val', '45000')  # 45% opacity
        print(f"  Slide 3 [{i}]: Added 45% opacity to green overlay")

print(f"\nSlide 3: Fixed {len(green_overlay_indices)} green overlays")

# ============================================================
# FIX 2: Slide 3 - Fix common component connections
# ============================================================
# Overlay [149] covers 06.공통 section, extend to include 07.포인트
shape149 = shapes3[149]
sp149 = shape149._element
spPr149 = sp149.find('p:spPr', NS)
if spPr149 is None:
    spPr149 = sp149.find('.//p:spPr', NS)
xfrm149 = spPr149.find('a:xfrm', NS)
off149 = xfrm149.find('a:off', NS)
ext149 = xfrm149.find('a:ext', NS)

# Extend to cover 공통+푸터+포인트 (y=1901952 to y=2642616)
new_height = 2642616 - 1901952
ext149.set('cy', str(new_height))
print(f"Slide 3: Extended overlay [149] to h={new_height}")

# Clone overlay for 08.관리자관리 section
admin_y = 2642616
admin_h = 3593592 - 2642616 + 237744

sp149_clone = copy.deepcopy(sp149)
xfrm_new = sp149_clone.find('p:spPr', NS)
if xfrm_new is None:
    xfrm_new = sp149_clone.find('.//p:spPr', NS)
xfrm_new = xfrm_new.find('a:xfrm', NS)
off_new = xfrm_new.find('a:off', NS)
ext_new = xfrm_new.find('a:ext', NS)
off_new.set('y', str(admin_y))
ext_new.set('cy', str(admin_h))

# Update shape ID
all_ids = []
for s in shapes3:
    cNvPr = s._element.find('.//p:cNvPr', NS)
    if cNvPr is not None:
        try:
            all_ids.append(int(cNvPr.get('id', '0')))
        except:
            pass
new_id = str(max(all_ids) + 10)
cNvPr_new = sp149_clone.find('.//p:cNvPr', NS)
cNvPr_new.set('id', new_id)
cNvPr_new.set('name', f'Shape {new_id}')

spTree3 = slide3._element.find('.//p:spTree', NS)
spTree3.append(sp149_clone)
print(f"Slide 3: Added overlay for 관리자관리 (y={admin_y}, h={admin_h})")

# ============================================================
# FIX 3: Slides 5-15 - Fix tag circles (green fill, white text, ellipse)
# ============================================================
def fix_tag_circle(sp):
    """Ensure tag circle has green fill, white bold text, ellipse shape"""
    spPr = sp.find('p:spPr', NS)
    if spPr is None:
        spPr = sp.find('.//p:spPr', NS)
    if spPr is None:
        return

    # Set green fill
    fill = spPr.find('a:solidFill', NS)
    if fill is not None:
        spPr.remove(fill)
    # Insert solidFill before a:ln or at start
    ln = spPr.find('a:ln', NS)
    fill = etree.Element(qn('a:solidFill'))
    clr = etree.SubElement(fill, qn('a:srgbClr'))
    clr.set('val', '3A9D6E')
    if ln is not None:
        ln.addprevious(fill)
    else:
        geom = spPr.find('a:prstGeom', NS)
        if geom is not None:
            geom.addprevious(fill)
        else:
            spPr.append(fill)

    # Set ellipse geometry
    geom = spPr.find('a:prstGeom', NS)
    if geom is not None:
        geom.set('prst', 'ellipse')

    # Remove outline
    ln = spPr.find('a:ln', NS)
    if ln is not None:
        spPr.remove(ln)
    ln_new = etree.SubElement(spPr, qn('a:ln'))

    # Set white bold text
    for rpr in sp.findall('.//a:rPr', NS):
        rpr.set('b', '1')
        for sf in rpr.findall('a:solidFill', NS):
            rpr.remove(sf)
        sf = etree.SubElement(rpr, qn('a:solidFill'))
        clr2 = etree.SubElement(sf, qn('a:srgbClr'))
        clr2.set('val', 'FFFFFF')

    # Center text
    for bodyPr in sp.findall('.//a:bodyPr', NS):
        bodyPr.set('anchor', 'ctr')
    for pPr in sp.findall('.//a:pPr', NS):
        pPr.set('algn', 'ctr')


for slide_idx in range(4, 15):  # Slides 5-15
    slide = prs.slides[slide_idx]
    fixed_count = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if len(text) == 1 and text in 'ABCDEFG':
                fix_tag_circle(shape._element)
                fixed_count += 1
    if fixed_count > 0:
        print(f"Slide {slide_idx+1}: Fixed {fixed_count} tag circles")

# ============================================================
# FIX 4: Style description labels (A., B., etc.) - green bold
# ============================================================
for slide_idx in range(4, 15):
    slide = prs.slides[slide_idx]
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if len(text) == 2 and text[0] in 'ABCDEFG' and text[1] == '.':
                for rpr in shape._element.findall('.//a:rPr', NS):
                    rpr.set('b', '1')
                    rpr.set('sz', '1200')
                    for sf in rpr.findall('a:solidFill', NS):
                        rpr.remove(sf)
                    sf = etree.SubElement(rpr, qn('a:solidFill'))
                    clr = etree.SubElement(sf, qn('a:srgbClr'))
                    clr.set('val', '3A9D6E')

print("Slides 5-15: Styled all A./B./C. labels")

# ============================================================
# FIX 5: Style description titles (bold, after A./B. labels)
# ============================================================
for slide_idx in range(4, 15):
    slide = prs.slides[slide_idx]
    shapes_list = list(slide.shapes)
    for i, shape in enumerate(shapes_list):
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            # Title shapes right after the label (same y, offset x)
            if len(text) == 2 and text[0] in 'ABCDEFG' and text[1] == '.':
                # Find the next shape at similar y position (title)
                label_y = shape.top
                for j in range(i+1, min(i+3, len(shapes_list))):
                    next_shape = shapes_list[j]
                    if next_shape.has_text_frame:
                        next_text = next_shape.text_frame.text.strip()
                        if abs(next_shape.top - label_y) < 50000 and len(next_text) > 2:
                            # This is the title - make bold
                            for rpr in next_shape._element.findall('.//a:rPr', NS):
                                rpr.set('b', '1')
                            break

print("Slides 5-15: Bolded description titles")

# ============================================================
# SAVE
# ============================================================
output_path = "김윤찬_TRY-CATCH_화면_포트폴리오_v7.pptx"
prs.save(output_path)
print(f"\nSaved to: {output_path}")
