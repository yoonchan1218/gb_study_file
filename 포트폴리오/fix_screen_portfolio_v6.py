from __future__ import annotations

import shutil
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.util import Pt


BASE_DIR = Path(r"C:\Users\pigch\Desktop\gb_0090_kyc\포트폴리오")
PPT_PATH = BASE_DIR / "김윤찬_TRY-CATCH_화면_포트폴리오_v6.pptx"
BACKUP_PATH = BASE_DIR / "김윤찬_TRY-CATCH_화면_포트폴리오_v6_before_final_fix_backup.pptx"
GENERATED_DIR = BASE_DIR / "_generated_screen_slides"
SCREENSHOT_DIR = BASE_DIR / "screenshots"

SLIDE_WIDTH = 9144000
SLIDE_HEIGHT = 5143500
PX_WIDTH = 1600
PX_HEIGHT = 900
SX = PX_WIDTH / SLIDE_WIDTH
SY = PX_HEIGHT / SLIDE_HEIGHT

GREEN = RGBColor(0x00, 0x6E, 0x3F)
SOFT_GREEN = RGBColor(0xDB, 0xEE, 0xE0)
BLACK = (40, 40, 40)
WHITE = (255, 255, 255)
LIGHT_BORDER = (90, 90, 90)

CATEGORY_MODAL_ASSET = SCREENSHOT_DIR / "v4-qna-category-modal.png"
PROFILE_ZOOM_ASSET = SCREENSHOT_DIR / "v4-mypage-profile-zoom.png"
EXPERIENCE_FILTER_ZOOM_ASSET = SCREENSHOT_DIR / "v4-experience-filter-zoom.png"
UNSUB_ERROR_ASSET = SCREENSHOT_DIR / "v4-unsub-error-state.png"


def emu_x(value: int) -> int:
    return round(value * SX)


def emu_y(value: int) -> int:
    return round(value * SY)


def ensure_backup() -> None:
    if not BACKUP_PATH.exists():
        shutil.copy2(PPT_PATH, BACKUP_PATH)


def remove_slide(prs: Presentation, index: int) -> None:
    slide_id = prs.slides._sldIdLst[index]
    rel_id = slide_id.rId
    prs.part.drop_rel(rel_id)
    prs.slides._sldIdLst.remove(slide_id)


def trim_to_base_slides(prs: Presentation, keep_count: int = 15) -> None:
    while len(prs.slides) > keep_count:
        remove_slide(prs, len(prs.slides) - 1)


def delete_shapes(slide, indices: list[int]) -> None:
    for idx in sorted(indices, reverse=True):
        shape = slide.shapes[idx]._element
        shape.getparent().remove(shape)


def move_pair(slide, shape_idx: int, text_idx: int, left: int, top: int) -> None:
    slide.shapes[shape_idx].left = left
    slide.shapes[shape_idx].top = top
    slide.shapes[text_idx].left = left
    slide.shapes[text_idx].top = top


def add_connector(slide, begin_x: int, begin_y: int, end_x: int, end_y: int) -> None:
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.ELBOW, begin_x, begin_y, end_x, end_y
    )
    connector.line.color.rgb = GREEN
    connector.line.width = Pt(1.3)


def label_anchor(slide, background_idx: int) -> tuple[int, int]:
    shape = slide.shapes[background_idx]
    return shape.left, shape.top + shape.height // 2


def asset_or_default(preferred: Path, fallback_name: str) -> Path:
    fallback = SCREENSHOT_DIR / fallback_name
    return preferred if preferred.exists() else fallback


def apply_soft_highlight(shape) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = SOFT_GREEN
    shape.fill.transparency = 0.45
    shape.line.color.rgb = GREEN
    shape.line.width = Pt(1.1)


def replace_picture_in_frame(slide, frame_idx: int, image_path: Path, padding: int = 9144) -> None:
    if not image_path.exists():
        return

    frame = slide.shapes[frame_idx]
    picture_indices: list[int] = []
    for idx, shape in enumerate(slide.shapes):
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        if (
            shape.left >= frame.left
            and shape.top >= frame.top
            and shape.left + shape.width <= frame.left + frame.width
            and shape.top + shape.height <= frame.top + frame.height
        ):
            picture_indices.append(idx)

    if picture_indices:
        delete_shapes(slide, picture_indices)

    inset = min(padding, frame.width // 8, frame.height // 8)
    slide.shapes.add_picture(
        str(image_path),
        frame.left + inset,
        frame.top + inset,
        width=frame.width - inset * 2,
        height=frame.height - inset * 2,
    )


def build_cropped_asset(
    source_path: Path,
    output_path: Path,
    output_size: tuple[int, int],
    crop_box: tuple[int, int, int, int] | None = None,
) -> None:
    if not source_path.exists():
        return

    image = Image.open(source_path).convert("RGB")
    if crop_box is not None:
        image = image.crop(crop_box)
    fitted = ImageOps.fit(image, output_size, method=Image.Resampling.LANCZOS)
    output_path.parent.mkdir(exist_ok=True)
    fitted.save(output_path)


def refresh_supporting_assets() -> None:
    build_cropped_asset(
        SCREENSHOT_DIR / "mypage-main-full.png",
        PROFILE_ZOOM_ASSET,
        (1000, 700),
        (250, 90, 1320, 840),
    )
    build_cropped_asset(
        SCREENSHOT_DIR / "mypage-experience-full.png",
        EXPERIENCE_FILTER_ZOOM_ASSET,
        (1600, 500),
        (470, 80, 1820, 500),
    )


def rebuild_connectors(
    slide,
    line_indices: list[int],
    background_indices: list[int],
    *,
    target_overrides: dict[int, tuple[int, int]] | None = None,
    pair_moves: dict[int, tuple[int, int, int]] | None = None,
) -> None:
    target_overrides = target_overrides or {}
    pair_moves = pair_moves or {}

    targets: dict[int, tuple[int, int]] = {}
    anchors: dict[int, tuple[int, int]] = {}
    for line_idx, background_idx in zip(line_indices, background_indices):
        line = slide.shapes[line_idx]
        targets[background_idx] = target_overrides.get(
            background_idx, (line.left, line.top)
        )

    for background_idx, (shape_idx, text_idx, top) in pair_moves.items():
        left = slide.shapes[shape_idx].left
        move_pair(slide, shape_idx, text_idx, left, top)

    for background_idx in background_indices:
        anchors[background_idx] = label_anchor(slide, background_idx)

    delete_shapes(slide, line_indices)

    for background_idx in background_indices:
        lx, ly = anchors[background_idx]
        tx, ty = targets[background_idx]
        add_connector(slide, tx, ty, lx, ly)


def fix_slide_3(prs: Presentation) -> None:
    slide = prs.slides[2]
    slide.shapes[77].width = 4974336
    for idx in [77, 148, 149]:
        apply_soft_highlight(slide.shapes[idx])


def fix_slide_4(prs: Presentation) -> None:
    slide = prs.slides[3]

    footer_top = 4562352
    point_top = 4821408
    move_pair(slide, 135, 136, 2286000, footer_top)
    move_pair(slide, 138, 139, 3383280, point_top)

    common = slide.shapes[129]
    header = slide.shapes[132]
    footer = slide.shapes[135]
    point = slide.shapes[138]

    common_right = common.left + common.width
    common_center_y = common.top + common.height // 2
    header_anchor = (header.left, header.top + header.height // 2)
    footer_anchor = (footer.left, footer.top + footer.height // 2)
    point_anchor = (point.left, point.top + point.height // 2)

    delete_shapes(slide, [131, 134, 137])

    add_connector(
        slide,
        common_right,
        common_center_y,
        header_anchor[0],
        header_anchor[1],
    )
    add_connector(
        slide,
        common_right,
        common_center_y,
        footer_anchor[0],
        footer_anchor[1],
    )
    add_connector(
        slide,
        common_right,
        common_center_y,
        point_anchor[0],
        point_anchor[1],
    )


def fix_slide_5(prs: Presentation) -> None:
    slide = prs.slides[4]
    rebuild_connectors(
        slide,
        [5, 8, 11, 14, 17, 20, 23],
        [6, 9, 12, 15, 18, 21, 24],
        target_overrides={
            6: (2797760, 848106),
            9: (640080, 1024128),
            12: (2171736, 1115568),
            15: (2171736, 1815084),
            18: (640080, 1767840),
            21: (1783080, 3870960),
            24: (2171736, 4572000),
        },
        pair_moves={
            15: (15, 16, 1554480),
            18: (18, 19, 2011680),
            21: (21, 22, 3764280),
        },
    )


def fix_slide_6(prs: Presentation) -> None:
    slide = prs.slides[5]
    rebuild_connectors(
        slide,
        [5, 8, 11, 14, 17],
        [6, 9, 12, 15, 18],
        target_overrides={
            6: (910000, 1440000),
            9: (3567113, 1700784),
            12: (3567113, 2226564),
            15: (2881312, 3131820),
            18: (3112034, 4078224),
        },
    )
    replace_picture_in_frame(
        slide,
        33,
        asset_or_default(CATEGORY_MODAL_ASSET, "detail-qna-write-category.png"),
    )


def fix_slide_7(prs: Presentation) -> None:
    slide = prs.slides[6]
    rebuild_connectors(
        slide,
        [5, 8, 11, 14, 17],
        [6, 9, 12, 15, 18],
        target_overrides={
            6: (1051560, 1070352),
            9: (3108960, 1221295),
            12: (3108960, 2236089),
            15: (3108960, 2671001),
            18: (3108960, 3685794),
        },
        pair_moves={
            6: (6, 7, 942336),
            9: (9, 10, 1234440),
        },
    )


def fix_slide_8(prs: Presentation) -> None:
    slide = prs.slides[7]
    rebuild_connectors(
        slide,
        [5, 8, 11, 14, 17],
        [6, 9, 12, 15, 18],
        target_overrides={
            6: (2520000, 1033272),
            9: (2720000, 1609344),
            12: (3020000, 2843784),
            15: (2940000, 3255264),
            18: (2860000, 4283964),
        },
    )


def fix_slide_9(prs: Presentation) -> None:
    slide = prs.slides[8]
    rebuild_connectors(
        slide,
        [5, 8, 11, 14],
        [6, 9, 12, 15],
        target_overrides={
            6: (3108960, 2535174),
            9: (3108960, 2792349),
            12: (3108960, 2998089),
            15: (3108960, 3178112),
        },
    )


def fix_slide_10(prs: Presentation) -> None:
    slide = prs.slides[9]
    rebuild_connectors(
        slide,
        [5, 8, 11, 14, 17, 20],
        [6, 9, 12, 15, 18, 21],
        target_overrides={
            6: (600000, 860000),
            9: (1500000, 1100000),
            12: (2740000, 1460000),
            15: (2820000, 1960000),
            18: (2780000, 2470000),
            21: (2050000, 3370000),
        },
        pair_moves={
            6: (6, 7, 731520),
            9: (9, 10, 1005840),
            12: (12, 13, 1307592),
        },
    )
    replace_picture_in_frame(
        slide,
        38,
        asset_or_default(PROFILE_ZOOM_ASSET, "detail-mypage-profile.png"),
    )


def fix_slide_11(prs: Presentation) -> None:
    slide = prs.slides[10]
    rebuild_connectors(
        slide,
        [5, 8, 11, 14, 17],
        [6, 9, 12, 15, 18],
        target_overrides={
            6: (2194560, 812101),
            9: (2500000, 992124),
            12: (2500000, 1249299),
            15: (2500000, 1557909),
            18: (2423160, 1866519),
        },
    )


def fix_slide_12(prs: Presentation) -> None:
    slide = prs.slides[11]
    rebuild_connectors(
        slide,
        [5, 8, 11, 14],
        [6, 9, 12, 15],
        target_overrides={
            6: (640080, 1306925),
            9: (2060000, 1140000),
            12: (3600000, 1140000),
            15: (2423160, 1515142),
        },
        pair_moves={
            6: (6, 7, 896112),
            9: (9, 10, 1115568),
            12: (12, 13, 1371600),
            15: (15, 16, 1636776),
        },
    )


def fix_slide_13(prs: Presentation) -> None:
    slide = prs.slides[12]
    rebuild_connectors(
        slide,
        [5, 8, 11, 14, 17],
        [6, 9, 12, 15, 18],
        target_overrides={
            6: (3108960, 1098614),
            9: (3108960, 1449872),
            12: (3108960, 2152388),
            15: (3108960, 2737818),
            18: (3108960, 3518392),
        },
    )
    replace_picture_in_frame(
        slide,
        33,
        asset_or_default(EXPERIENCE_FILTER_ZOOM_ASSET, "fix-exp-filter.png"),
    )


def fix_slide_14(prs: Presentation) -> None:
    slide = prs.slides[13]
    rebuild_connectors(
        slide,
        [5, 8, 11, 14],
        [6, 9, 12, 15],
        target_overrides={
            6: (3108960, 850583),
            9: (3108960, 1299972),
            12: (3108960, 1813560),
            15: (3108960, 2134553),
        },
    )
    replace_picture_in_frame(
        slide,
        28,
        asset_or_default(UNSUB_ERROR_ASSET, "fix-unsub-error.png"),
    )


def load_font(name: str, size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        Path(r"C:\Windows\Fonts") / name,
        Path(r"C:\Windows\Fonts") / "malgun.ttf",
    ]
    for candidate in candidates:
        if candidate.exists():
            return ImageFont.truetype(str(candidate), size=size)
    return ImageFont.load_default()


FONT_BOLD_34 = load_font("malgunbd.ttf", 34)
FONT_BOLD_28 = load_font("malgunbd.ttf", 28)
FONT_BOLD_26 = load_font("malgunbd.ttf", 26)
FONT_BOLD_24 = load_font("malgunbd.ttf", 24)
FONT_BOLD_22 = load_font("malgunbd.ttf", 22)
FONT_BOLD_20 = load_font("malgunbd.ttf", 20)
FONT_REG_20 = load_font("malgun.ttf", 20)
FONT_REG_18 = load_font("malgun.ttf", 18)
FONT_REG_16 = load_font("malgun.ttf", 16)


def fit_image(path: Path, width: int, height: int) -> Image.Image:
    image = Image.open(path).convert("RGB")
    contained = ImageOps.contain(image, (width, height), Image.Resampling.LANCZOS)
    canvas = Image.new("RGB", (width, height), WHITE)
    offset_x = (width - contained.width) // 2
    offset_y = (height - contained.height) // 2
    canvas.paste(contained, (offset_x, offset_y))
    return canvas


def draw_box(draw: ImageDraw.ImageDraw, rect: tuple[int, int, int, int], width: int = 2) -> None:
    draw.rectangle(rect, outline=LIGHT_BORDER, width=width)


def wrap_text(draw: ImageDraw.ImageDraw, text: str, font, max_width: int) -> str:
    words = text.split()
    if not words:
        return text

    lines: list[str] = []
    current = words[0]
    for word in words[1:]:
        candidate = f"{current} {word}"
        if draw.textbbox((0, 0), candidate, font=font)[2] <= max_width:
            current = candidate
        else:
            lines.append(current)
            current = word
    lines.append(current)
    return "\n".join(lines)


def draw_text(
    draw: ImageDraw.ImageDraw,
    rect: tuple[int, int, int, int],
    text: str,
    font,
    *,
    bold: bool = False,
    align: str = "left",
    line_spacing: int = 6,
) -> None:
    x, y, w, h = rect
    wrapped = wrap_text(draw, text, font, max(10, w - 16))
    lines = wrapped.split("\n")
    cursor_y = y + 8
    for line in lines:
        bbox = draw.textbbox((0, 0), line, font=font)
        line_width = bbox[2] - bbox[0]
        if align == "center":
            line_x = x + (w - line_width) // 2
        else:
            line_x = x + 8
        draw.text((line_x, cursor_y), line, fill=BLACK, font=font)
        cursor_y += (bbox[3] - bbox[1]) + line_spacing


def paste_frame(
    canvas: Image.Image,
    draw: ImageDraw.ImageDraw,
    title_rect_emu: tuple[int, int, int, int],
    frame_rect_emu: tuple[int, int, int, int],
    image_path: Path,
    title: str,
) -> None:
    title_rect = (
        emu_x(title_rect_emu[0]),
        emu_y(title_rect_emu[1]),
        emu_x(title_rect_emu[2]),
        emu_y(title_rect_emu[3]),
    )
    draw_text(draw, title_rect, title, FONT_BOLD_22)

    fx = emu_x(frame_rect_emu[0])
    fy = emu_y(frame_rect_emu[1])
    fw = emu_x(frame_rect_emu[2])
    fh = emu_y(frame_rect_emu[3])
    draw_box(draw, (fx, fy, fx + fw, fy + fh))
    fitted = fit_image(image_path, fw - 8, fh - 8)
    canvas.paste(fitted, (fx + 4, fy + 4))


def render_detail_slide(
    slide_no: int,
    subtitle: str,
    right_title: str,
    bullets: list[tuple[str, str, str]],
    top_left: tuple[str, str],
    top_right: tuple[str, str],
    bottom_left: tuple[str, str],
    mini: tuple[str, str],
) -> Path:
    GENERATED_DIR.mkdir(exist_ok=True)

    image = Image.new("RGB", (PX_WIDTH, PX_HEIGHT), WHITE)
    draw = ImageDraw.Draw(image)

    draw.rectangle((0, 0, emu_x(54864), emu_y(457200)), fill=(0, 0, 0))

    title_box = (emu_x(182880), emu_y(45720), emu_x(8229600), emu_y(320040))
    subtitle_box = (emu_x(182880), emu_y(384048), emu_x(4572000), emu_y(201168))
    draw_box(draw, (title_box[0], title_box[1], title_box[0] + title_box[2], title_box[1] + title_box[3]))
    draw_box(draw, (subtitle_box[0], subtitle_box[1], subtitle_box[0] + subtitle_box[2], subtitle_box[1] + subtitle_box[3]))
    draw_text(
        draw,
        (title_box[0], title_box[1], title_box[2], title_box[3]),
        "TRY-CATCH    화면 내 영역 설명",
        FONT_BOLD_28,
    )
    draw_text(
        draw,
        (subtitle_box[0], subtitle_box[1], subtitle_box[2], subtitle_box[3]),
        subtitle,
        FONT_BOLD_26,
    )

    paste_frame(
        image,
        draw,
        (137160, 640080, 4206240, 146304),
        (128016, 795528, 3269488, 1847088),
        SCREENSHOT_DIR / top_left[1],
        top_left[0],
    )
    paste_frame(
        image,
        draw,
        (4663440, 640080, 4297680, 146304),
        (4654296, 795528, 3269488, 1847088),
        SCREENSHOT_DIR / top_right[1],
        top_right[0],
    )
    paste_frame(
        image,
        draw,
        (137160, 2761488, 4206240, 146304),
        (128016, 2916936, 3594608, 2029968),
        SCREENSHOT_DIR / bottom_left[1],
        bottom_left[0],
    )
    paste_frame(
        image,
        draw,
        (137160, 4544568, 2286000, 146304),
        (128016, 4700016, 424688, 384048),
        SCREENSHOT_DIR / mini[1],
        mini[0],
    )

    right_title_rect = (emu_x(4663440), emu_y(2743200), emu_x(4297680), emu_y(182880))
    draw_box(
        draw,
        (
            right_title_rect[0],
            right_title_rect[1],
            right_title_rect[0] + right_title_rect[2],
            right_title_rect[1] + right_title_rect[3],
        ),
    )
    draw_text(draw, right_title_rect, right_title, FONT_BOLD_24)

    y_positions = [3017520, 3401568, 3785616, 4169664, 4553712]
    for idx, (label, title, desc) in enumerate(bullets):
        row_y = y_positions[idx]
        label_rect = (emu_x(4663440), emu_y(row_y), emu_x(228600), emu_y(182880))
        title_rect = (emu_x(4892040), emu_y(row_y), emu_x(4069080), emu_y(128016))
        desc_rect = (emu_x(4892040), emu_y(row_y + 128016), emu_x(4069080), emu_y(237744))

        draw_box(draw, (label_rect[0], label_rect[1], label_rect[0] + label_rect[2], label_rect[1] + label_rect[3]))
        draw_box(draw, (title_rect[0], title_rect[1], title_rect[0] + title_rect[2], title_rect[1] + title_rect[3]))
        draw_box(draw, (desc_rect[0], desc_rect[1], desc_rect[0] + desc_rect[2], desc_rect[1] + desc_rect[3]))

        draw_text(draw, label_rect, label, FONT_BOLD_22)
        draw_text(draw, title_rect, title, FONT_BOLD_20)
        draw_text(draw, desc_rect, desc, FONT_REG_18, line_spacing=4)

    page_box = (emu_x(8686800), emu_y(4846320), emu_x(274320), emu_y(182880))
    draw_box(draw, (page_box[0], page_box[1], page_box[0] + page_box[2], page_box[1] + page_box[3]))
    draw_text(draw, page_box, str(slide_no), FONT_BOLD_24, align="center")

    output_path = GENERATED_DIR / f"slide_{slide_no:02d}.png"
    image.save(output_path)
    return output_path


def append_generated_slides(prs: Presentation) -> None:
    slide_images = [
        render_detail_slide(
            16,
            "〈Q&A 작성/수정 세부 인터랙션〉",
            "Q&A 작성/수정 기능 단위 설명",
            [
                ("A.", "카테고리 드롭다운", "직무/기업/동문 대상 선택 레이어"),
                ("B.", "질문 대상 토글", "체크 상태 즉시 반영, 선택값 유지"),
                ("C.", "수정 화면 프리로드", "기존 제목/본문/카테고리 그대로 재편집"),
                ("D.", "첨부파일/에디터", "이미지·파일 첨부, 기존 첨부 관리"),
                ("E.", "예외 메시지", "필수값 누락 시 alert/경고 문구 노출"),
            ],
            ("카테고리 선택 드롭다운", asset_or_default(CATEGORY_MODAL_ASSET, "detail-qna-write-category.png").name),
            ("질문 대상 선택 / 체크 상태", "fix-qna-write-checked.png"),
            ("수정 화면 + 첨부 관리", "qna-update.png"),
            ("예외 alert", "detail-qna-write-alert.png"),
        ),
        render_detail_slide(
            17,
            "〈체험 지원 첨부파일 인터랙션〉",
            "체험 지원/첨부 기능 단위 설명",
            [
                ("A.", "지원서 폼", "공고 상세에서 첨부파일 섹션 확인"),
                ("B.", "파일/URL 토글", "버튼 전환으로 업로드 방식 변경"),
                ("C.", "첨부 카테고리 선택", "증빙/포트폴리오 등 분류 선택"),
                ("D.", "업로드 결과 표시", "파일명·용량·삭제 버튼 즉시 반영"),
                ("E.", "지원 현황 상태", "마이페이지에서 상태 필터/진행 단계 확인"),
            ],
            ("지원 폼 - 첨부파일 섹션", "_training_program_file.png"),
            ("파일첨부 팝업 + 타입 토글", "training-program-attach-popup.png"),
            ("첨부 결과 / 파일 목록 반영", "training-program-attach-result.png"),
            ("상태 필터", asset_or_default(EXPERIENCE_FILTER_ZOOM_ASSET, "fix-exp-filter.png").name),
        ),
        render_detail_slide(
            18,
            "〈공통 상태 변화: 로그인 / 알림 / 검색〉",
            "공통 헤더 상태 변화 기능 설명",
            [
                ("A.", "비로그인 상태", "기업서비스 버튼 중심 GNB 진입"),
                ("B.", "비회원 드롭다운", "로그인 전 기업서비스/가입 유도"),
                ("C.", "로그인 후 프로필", "회원명 클릭 시 개인 메뉴 노출"),
                ("D.", "알림 아이콘/드롭다운", "새 알림 표시와 목록 진입"),
                ("E.", "검색 자동완성", "추천 키워드 + 자동완성 끄기 제공"),
            ],
            ("비로그인 - 기업서비스 드롭다운", "v3-header-dropdown-before.png"),
            ("로그인 후 - 프로필 드롭다운", "v3-header-profile-after.png"),
            ("알림 드롭다운", "header-notification.png"),
            ("검색 자동완성", "v3-search-dropdown.png"),
        ),
    ]

    blank_layout = prs.slide_layouts[0]
    for slide_image in slide_images:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(slide_image),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )


def main() -> None:
    ensure_backup()
    refresh_supporting_assets()
    prs = Presentation(str(PPT_PATH))
    trim_to_base_slides(prs)

    fix_slide_3(prs)
    fix_slide_4(prs)
    fix_slide_5(prs)
    fix_slide_6(prs)
    fix_slide_7(prs)
    fix_slide_8(prs)
    fix_slide_9(prs)
    fix_slide_10(prs)
    fix_slide_11(prs)
    fix_slide_12(prs)
    fix_slide_13(prs)
    fix_slide_14(prs)
    append_generated_slides(prs)

    prs.save(str(PPT_PATH))


if __name__ == "__main__":
    main()
