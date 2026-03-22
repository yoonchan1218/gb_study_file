"""
TRY-CATCH 포트폴리오 최종 수정
- python-pptx만 사용, unpack/pack 안 함
"""
import copy
from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}

prs = Presentation("김윤찬_TRY-CATCH_화면_포트폴리오_v6.pptx")

# ============================================================
# UTIL: Delete slide properly
# ============================================================
def delete_slide(prs, index):
    rId = prs.slides._sldIdLst[index].get(qn('r:id'))
    prs.part.drop_rel(rId)
    sldId = prs.slides._sldIdLst[index]
    prs.slides._sldIdLst.remove(sldId)

def move_slide(prs, old_index, new_index):
    slides = prs.slides._sldIdLst
    el = slides[old_index]
    slides.remove(el)
    if new_index >= len(slides):
        slides.append(el)
    else:
        slides.insert(new_index, el)

# ============================================================
# STEP 1: Slide 3 - Green overlay transparency
# ============================================================
print("=== STEP 1: Slide 3 green overlays ===")
slide3 = prs.slides[2]
count = 0
for shape in slide3.shapes:
    sp = shape._element
    for clr in sp.findall('.//' + qn('a:srgbClr')):
        if clr.get('val') == 'DBEEE0':
            for old in clr.findall(qn('a:alpha')):
                clr.remove(old)
            alpha = etree.SubElement(clr, qn('a:alpha'))
            alpha.set('val', '45000')
            count += 1
print(f"  Added transparency to {count} green fills")

# ============================================================
# STEP 2: Slides 5-15 - Remove connectors, fix tags
# ============================================================
print("\n=== STEP 2: Remove connectors, fix tags ===")
for slide_idx in range(4, 15):
    slide = prs.slides[slide_idx]
    spTree = slide._element.find('.//' + qn('p:spTree'))

    # Remove all cxnSp (connector) elements
    connectors = spTree.findall(qn('p:cxnSp'))
    for c in connectors:
        spTree.remove(c)

    # Find main image bounds for tag repositioning
    img_right = 0
    img_left = 0
    for shape in slide.shapes:
        try:
            _ = shape.image
            if shape.width > 2000000:  # main screenshot image
                img_right = shape.left + shape.width
                img_left = shape.left
                break
        except:
            pass

    TAG_SIZE = 146304  # ~0.16 inches
    TAG_INSET = 100000  # how far inside the image right edge

    # Collect tag bg+text pairs by position
    tag_fixed = 0
    for shape in list(slide.shapes):
        if not shape.has_text_frame:
            sp = shape._element
            spPr = sp.find(qn('p:spPr'))
            if spPr is None:
                continue
            fill = spPr.find(qn('a:solidFill'))
            if fill is None:
                continue
            clr = fill.find(qn('a:srgbClr'))
            if clr is not None and clr.get('val') == '3A9D6E':
                if shape.width == 256032:
                    # Move inside image + shrink + ellipse
                    new_x = img_right - TAG_SIZE - TAG_INSET
                    center_y = shape.top + shape.height // 2
                    shape.left = new_x
                    shape.top = center_y - TAG_SIZE // 2
                    shape.width = TAG_SIZE
                    shape.height = TAG_SIZE
                    geom = spPr.find(qn('a:prstGeom'))
                    if geom is not None:
                        geom.set('prst', 'ellipse')
                    ln = spPr.find(qn('a:ln'))
                    if ln is not None:
                        spPr.remove(ln)
                    etree.SubElement(spPr, qn('a:ln'))
                    tag_fixed += 1
            continue

        text = shape.text_frame.text.strip()

        # Single letter tag (A-G)
        if len(text) == 1 and text in 'ABCDEFG':
            sp = shape._element
            spPr = sp.find(qn('p:spPr'))
            if spPr is not None and shape.width == 256032:
                new_x = img_right - TAG_SIZE - TAG_INSET
                center_y = shape.top + shape.height // 2
                shape.left = new_x
                shape.top = center_y - TAG_SIZE // 2
                shape.width = TAG_SIZE
                shape.height = TAG_SIZE
                geom = spPr.find(qn('a:prstGeom'))
                if geom is not None:
                    geom.set('prst', 'ellipse')

            for rpr in sp.findall('.//' + qn('a:rPr')):
                rpr.set('sz', '700')
                if rpr.get('b') is None:
                    rpr.set('b', '1')
                for sf in rpr.findall(qn('a:solidFill')):
                    rpr.remove(sf)
                sf = etree.SubElement(rpr, qn('a:solidFill'))
                c = etree.SubElement(sf, qn('a:srgbClr'))
                c.set('val', 'FFFFFF')
            for bp in sp.findall('.//' + qn('a:bodyPr')):
                bp.set('anchor', 'ctr')
                bp.set('lIns', '0')
                bp.set('tIns', '0')
                bp.set('rIns', '0')
                bp.set('bIns', '0')
            for pp in sp.findall('.//' + qn('a:pPr')):
                pp.set('algn', 'ctr')
            tag_fixed += 1

        # Label (A., B., etc.)
        elif len(text) == 2 and text[0] in 'ABCDEFG' and text[1] == '.':
            for rpr in shape._element.findall('.//' + qn('a:rPr')):
                if rpr.get('b') is None:
                    rpr.set('b', '1')
                for sf in rpr.findall(qn('a:solidFill')):
                    rpr.remove(sf)
                sf = etree.SubElement(rpr, qn('a:solidFill'))
                c = etree.SubElement(sf, qn('a:srgbClr'))
                c.set('val', '3A9D6E')

    if len(connectors) > 0 or tag_fixed > 0:
        print(f"  Slide {slide_idx+1}: {len(connectors)} connectors removed, {tag_fixed} tags fixed")

# ============================================================
# STEP 3: Fix image-only slides (16, 18) - add header like others
# ============================================================
print("\n=== STEP 3: Fix image-only slides ===")

def add_header_to_image_slide(slide, subtitle_text):
    """Add standard header bar + subtitle to an image-only slide"""
    spTree = slide._element.find('.//' + qn('p:spTree'))

    # Resize existing image to leave room for header
    for shape in slide.shapes:
        try:
            _ = shape.image
            # Move image down and shrink
            shape.left = 0
            shape.top = 457200  # below header
            shape.width = 9144000
            shape.height = 4686300  # rest of slide
        except:
            pass

    # Add dark accent bar (left side thin stripe)
    bar_sp = etree.SubElement(spTree, qn('p:sp'))
    nvSpPr = etree.SubElement(bar_sp, qn('p:nvSpPr'))
    cNvPr = etree.SubElement(nvSpPr, qn('p:cNvPr'))
    cNvPr.set('id', '200')
    cNvPr.set('name', 'HeaderBar')
    etree.SubElement(nvSpPr, qn('p:cNvSpPr'))
    etree.SubElement(nvSpPr, qn('p:nvPr'))
    spPr = etree.SubElement(bar_sp, qn('p:spPr'))
    xfrm = etree.SubElement(spPr, qn('a:xfrm'))
    off = etree.SubElement(xfrm, qn('a:off'))
    off.set('x', '0')
    off.set('y', '0')
    ext = etree.SubElement(xfrm, qn('a:ext'))
    ext.set('cx', '54864')
    ext.set('cy', '457200')
    geom = etree.SubElement(spPr, qn('a:prstGeom'))
    geom.set('prst', 'rect')
    etree.SubElement(geom, qn('a:avLst'))
    fill = etree.SubElement(spPr, qn('a:solidFill'))
    clr = etree.SubElement(fill, qn('a:srgbClr'))
    clr.set('val', '333333')
    etree.SubElement(spPr, qn('a:ln'))
    txBody = etree.SubElement(bar_sp, qn('p:txBody'))
    etree.SubElement(txBody, qn('a:bodyPr'))
    etree.SubElement(txBody, qn('a:p'))

    # Add title text
    title_sp = etree.SubElement(spTree, qn('p:sp'))
    nvSpPr2 = etree.SubElement(title_sp, qn('p:nvSpPr'))
    cNvPr2 = etree.SubElement(nvSpPr2, qn('p:cNvPr'))
    cNvPr2.set('id', '201')
    cNvPr2.set('name', 'TitleText')
    etree.SubElement(nvSpPr2, qn('p:cNvSpPr'))
    etree.SubElement(nvSpPr2, qn('p:nvPr'))
    spPr2 = etree.SubElement(title_sp, qn('p:spPr'))
    xfrm2 = etree.SubElement(spPr2, qn('a:xfrm'))
    off2 = etree.SubElement(xfrm2, qn('a:off'))
    off2.set('x', '182880')
    off2.set('y', '45720')
    ext2 = etree.SubElement(xfrm2, qn('a:ext'))
    ext2.set('cx', '8229600')
    ext2.set('cy', '320040')
    geom2 = etree.SubElement(spPr2, qn('a:prstGeom'))
    geom2.set('prst', 'rect')
    etree.SubElement(geom2, qn('a:avLst'))
    etree.SubElement(spPr2, qn('a:noFill'))
    etree.SubElement(spPr2, qn('a:ln'))
    txBody2 = etree.SubElement(title_sp, qn('p:txBody'))
    bodyPr2 = etree.SubElement(txBody2, qn('a:bodyPr'))
    bodyPr2.set('wrap', 'square')
    bodyPr2.set('lIns', '0')
    bodyPr2.set('tIns', '0')
    bodyPr2.set('rIns', '0')
    bodyPr2.set('bIns', '0')
    bodyPr2.set('rtlCol', '0')
    bodyPr2.set('anchor', 'ctr')
    etree.SubElement(txBody2, qn('a:lstStyle'))
    p = etree.SubElement(txBody2, qn('a:p'))
    pPr = etree.SubElement(p, qn('a:pPr'))
    pPr.set('indent', '0')
    pPr.set('marL', '0')
    etree.SubElement(pPr, qn('a:buNone'))
    r = etree.SubElement(p, qn('a:r'))
    rPr = etree.SubElement(r, qn('a:rPr'))
    rPr.set('lang', 'en-US')
    rPr.set('sz', '1400')
    rPr.set('b', '1')
    rPr.set('dirty', '0')
    sfill = etree.SubElement(rPr, qn('a:solidFill'))
    sclr = etree.SubElement(sfill, qn('a:srgbClr'))
    sclr.set('val', '333333')
    latin = etree.SubElement(rPr, qn('a:latin'))
    latin.set('typeface', '\ub9d1\uc740 \uace0\ub515')
    latin.set('pitchFamily', '34')
    latin.set('charset', '0')
    ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', '\ub9d1\uc740 \uace0\ub515')
    ea.set('pitchFamily', '34')
    ea.set('charset', '-122')
    t = etree.SubElement(r, qn('a:t'))
    t.text = 'TRY-CATCH    \ud654\uba74 \ub0b4 \uc601\uc5ed \uc124\uba85'

    # Add subtitle
    sub_sp = etree.SubElement(spTree, qn('p:sp'))
    nvSpPr3 = etree.SubElement(sub_sp, qn('p:nvSpPr'))
    cNvPr3 = etree.SubElement(nvSpPr3, qn('p:cNvPr'))
    cNvPr3.set('id', '202')
    cNvPr3.set('name', 'SubTitle')
    etree.SubElement(nvSpPr3, qn('p:cNvSpPr'))
    etree.SubElement(nvSpPr3, qn('p:nvPr'))
    spPr3 = etree.SubElement(sub_sp, qn('p:spPr'))
    xfrm3 = etree.SubElement(spPr3, qn('a:xfrm'))
    off3 = etree.SubElement(xfrm3, qn('a:off'))
    off3.set('x', '182880')
    off3.set('y', '384048')
    ext3 = etree.SubElement(xfrm3, qn('a:ext'))
    ext3.set('cx', '4572000')
    ext3.set('cy', '201168')
    geom3 = etree.SubElement(spPr3, qn('a:prstGeom'))
    geom3.set('prst', 'rect')
    etree.SubElement(geom3, qn('a:avLst'))
    etree.SubElement(spPr3, qn('a:noFill'))
    etree.SubElement(spPr3, qn('a:ln'))
    txBody3 = etree.SubElement(sub_sp, qn('p:txBody'))
    bodyPr3 = etree.SubElement(txBody3, qn('a:bodyPr'))
    bodyPr3.set('wrap', 'square')
    bodyPr3.set('lIns', '0')
    bodyPr3.set('tIns', '0')
    bodyPr3.set('rIns', '0')
    bodyPr3.set('bIns', '0')
    bodyPr3.set('anchor', 'ctr')
    etree.SubElement(txBody3, qn('a:lstStyle'))
    p3 = etree.SubElement(txBody3, qn('a:p'))
    pPr3 = etree.SubElement(p3, qn('a:pPr'))
    pPr3.set('indent', '0')
    pPr3.set('marL', '0')
    etree.SubElement(pPr3, qn('a:buNone'))
    r3 = etree.SubElement(p3, qn('a:r'))
    rPr3 = etree.SubElement(r3, qn('a:rPr'))
    rPr3.set('lang', 'en-US')
    rPr3.set('sz', '1000')
    rPr3.set('dirty', '0')
    sfill3 = etree.SubElement(rPr3, qn('a:solidFill'))
    sclr3 = etree.SubElement(sfill3, qn('a:srgbClr'))
    sclr3.set('val', '666666')
    latin3 = etree.SubElement(rPr3, qn('a:latin'))
    latin3.set('typeface', '\ub9d1\uc740 \uace0\ub515')
    latin3.set('pitchFamily', '34')
    latin3.set('charset', '0')
    ea3 = etree.SubElement(rPr3, qn('a:ea'))
    ea3.set('typeface', '\ub9d1\uc740 \uace0\ub515')
    ea3.set('pitchFamily', '34')
    ea3.set('charset', '-122')
    t3 = etree.SubElement(r3, qn('a:t'))
    t3.text = subtitle_text

# Slide 16 (0-indexed 15) = Q&A 작성/수정 세부 인터랙션
slide16 = prs.slides[15]
add_header_to_image_slide(slide16, '\u300cQ&A \uc791\uc131/\uc218\uc815 \uc138\ubd80 \uc778\ud130\ub799\uc158\u300d')
print("  Slide 16: Added header (Q&A 작성/수정 세부 인터랙션)")

# Slide 18 (0-indexed 17) = 공통 헤더 상태 변화
slide18 = prs.slides[17]
add_header_to_image_slide(slide18, '\u300c\uacf5\ud1b5 \uc0c1\ud0dc \ubcc0\ud654: \ub85c\uadf8\uc778 / \uc54c\ub9bc / \uac80\uc0c9\u300d')
print("  Slide 18: Added header (공통 상태 변화)")

# ============================================================
# STEP 4: Delete slide 17 (체험 지원 - not my work)
# ============================================================
print("\n=== STEP 4: Delete slide 17 ===")
delete_slide(prs, 16)  # 0-indexed 16 = slide 17
print("  Deleted slide 17 (체험 지원 첨부파일)")

# ============================================================
# STEP 5: Reorder slides
# ============================================================
print("\n=== STEP 5: Reorder slides ===")
# After deleting slide 17, current order (17 slides):
# 0: 표지
# 1: 기획배경
# 2: 기능리스트
# 3: 메뉴트리
# 4: Q&A 목록       (slide5)
# 5: Q&A 작성       (slide6)
# 6: Q&A 상세       (slide7)
# 7: Q&A 수정       (slide8)
# 8: 공통 푸터       (slide9)
# 9: 마이페이지 홈   (slide10)
# 10: 개인정보 수정  (slide11)
# 11: 알림 목록      (slide12)
# 12: 활동 내역      (slide13)
# 13: 회원 탈퇴      (slide14)
# 14: 공통 드롭다운  (slide15)
# 15: Q&A 세부       (slide16)  → move after Q&A 수정
# 16: 공통 헤더      (slide18)  → move after 공통 드롭다운

# Move slide16(Q&A세부) from index 15 to after index 7 (after Q&A수정)
move_slide(prs, 15, 8)
# Now: 0,1,2,3,4,5,6,7, Q&A세부, 공통푸터, 마이페이지...공통드롭다운, 공통헤더
# Current indexes:
# 8: Q&A 세부
# 9: 공통 푸터
# ...
# 14: 공통 드롭다운
# 15: 공통 헤더 (was at 16, now at 15 after the move)

# Move 공통 푸터(index 9) to after 회원탈퇴(index 13 after move)
# Current: 0,1,2,3, 4(Q&A목록),5,6,7, 8(Q&A세부), 9(푸터), 10(마이페이지),11,12,13(탈퇴), 14(드롭다운), 15(헤더)
# Move 푸터(9) to after 탈퇴(13) → insert at 14
move_slide(prs, 9, 14)
# Now: 0,1,2,3, 4,5,6,7,8(Q&A세부), 9(마이페이지),10,11,12(탈퇴), 13(푸터), 14(드롭다운), 15(헤더)

print("  Reordered: Q&A(5-9) → 마이페이지(10-14) → 공통(15-17)")

# ============================================================
# STEP 6: Update page numbers
# ============================================================
print("\n=== STEP 6: Update page numbers ===")
for idx, slide in enumerate(prs.slides):
    page_num = idx + 1
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text.isdigit() and shape.left > 8000000 and shape.top > 4500000:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip().isdigit():
                            run.text = str(page_num)

# ============================================================
# SAVE
# ============================================================
output = "김윤찬_TRY-CATCH_화면_포트폴리오_v7_final.pptx"
prs.save(output)
import os
print(f"\nSaved: {output} ({os.path.getsize(output):,} bytes)")
print(f"Total slides: {len(prs.slides)}")

# Print final order
print("\nFinal slide order:")
labels = [
    "표지", "기획 배경", "기능 리스트", "메뉴트리",
    "Q&A 목록", "Q&A 작성", "Q&A 상세", "Q&A 수정", "Q&A 세부 인터랙션",
    "마이페이지 홈", "개인정보 수정", "알림 목록", "활동 내역", "회원 탈퇴",
    "공통 푸터", "공통 드롭다운", "공통 헤더 상태"
]
for i, label in enumerate(labels):
    print(f"  {i+1:2d}. {label}")
