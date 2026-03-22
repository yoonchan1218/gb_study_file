"""
TRY-CATCH 포트폴리오 v6 -> v7 수정 스크립트

수정 사항:
1. 슬라이드 3: 초록 오버레이 투명도 추가 (텍스트 보이기)
2. 슬라이드 17 제거 (내 작업 아님)
3. 슬라이드 순서 정리: Q&A(5-8,16) → 마이페이지(10-14) → 공통(9,15,18)
4. 슬라이드 5~15: connector 선 제거, 태그 작은 원형으로 변경
"""

import copy
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}

prs = Presentation("김윤찬_TRY-CATCH_화면_포트폴리오_v6.pptx")

# ============================================================
# STEP 1: Slide 3 - Green overlay transparency
# ============================================================
print("=== STEP 1: Slide 3 green overlay transparency ===")
slide3 = prs.slides[2]
for i, shape in enumerate(slide3.shapes):
    sp = shape._element
    spPr = sp.find('p:spPr', NS)
    if spPr is None:
        continue
    fill = spPr.find('a:solidFill', NS)
    if fill is None:
        continue
    clr = fill.find('a:srgbClr', NS)
    if clr is not None and clr.get('val') == 'DBEEE0':
        for old_alpha in clr.findall('a:alpha', NS):
            clr.remove(old_alpha)
        alpha = etree.SubElement(clr, qn('a:alpha'))
        alpha.set('val', '45000')
        print(f"  [{i}] Added 45% opacity")

# ============================================================
# STEP 2: Remove connector lines + fix tags (slides 5-15)
# ============================================================
print("\n=== STEP 2: Remove connectors, fix tag circles ===")

TAG_SIZE = 164592  # ~0.18 inches, smaller circles
TAG_OFFSET = 10000  # small inward offset

for slide_idx in range(4, 15):  # slides 5-15 (0-indexed 4-14)
    slide = prs.slides[slide_idx]
    spTree = slide._element.find('.//p:spTree', NS)

    shapes_to_remove = []
    tag_bg_shapes = []  # green background shapes for tags
    tag_text_shapes = []  # text shapes with A, B, C...

    for i, shape in enumerate(slide.shapes):
        sp = shape._element
        spPr = sp.find('p:spPr', NS)
        if spPr is None:
            continue

        geom = spPr.find('a:prstGeom', NS)

        # Remove bentConnector lines
        if geom is not None and 'Connector' in (geom.get('prst') or ''):
            shapes_to_remove.append(sp)
            continue
        if geom is not None and 'connector' in (geom.get('prst') or '').lower():
            shapes_to_remove.append(sp)
            continue

        # Also remove zero-width or zero-height line shapes that aren't tags
        w = shape.width
        h = shape.height
        if (w == 0 or h == 0) and not shape.has_text_frame:
            shapes_to_remove.append(sp)
            continue

        # Identify tag elements
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if len(text) == 1 and text in 'ABCDEFG':
                tag_text_shapes.append((shape, text))
            # Check for tag background (green rect at same position)
            fill = spPr.find('a:solidFill/a:srgbClr', NS)
            if fill is not None and fill.get('val') == '3A9D6E' and not shape.has_text_frame:
                tag_bg_shapes.append(shape)
        else:
            fill = spPr.find('a:solidFill/a:srgbClr', NS)
            if fill is not None and fill.get('val') == '3A9D6E':
                # This is a tag background shape
                if w == 256032 and h == 256032:
                    tag_bg_shapes.append(shape)

    # Remove connectors
    removed = 0
    for sp in shapes_to_remove:
        parent = sp.getparent()
        if parent is not None:
            parent.remove(sp)
            removed += 1

    # Fix tag background shapes: make smaller ellipse
    for shape in tag_bg_shapes:
        sp = shape._element
        spPr = sp.find('p:spPr', NS)
        xfrm = spPr.find('a:xfrm', NS)
        off = xfrm.find('a:off', NS)
        ext = xfrm.find('a:ext', NS)

        old_x = int(off.get('x'))
        old_y = int(off.get('y'))
        old_w = int(ext.get('cx'))
        old_h = int(ext.get('cy'))

        # Center the smaller circle at the same center point
        center_x = old_x + old_w // 2
        center_y = old_y + old_h // 2
        new_x = center_x - TAG_SIZE // 2
        new_y = center_y - TAG_SIZE // 2

        off.set('x', str(new_x))
        off.set('y', str(new_y))
        ext.set('cx', str(TAG_SIZE))
        ext.set('cy', str(TAG_SIZE))

        # Change to ellipse
        geom = spPr.find('a:prstGeom', NS)
        if geom is not None:
            geom.set('prst', 'ellipse')

        # Remove border
        ln = spPr.find('a:ln', NS)
        if ln is not None:
            spPr.remove(ln)
        ln_new = etree.SubElement(spPr, qn('a:ln'))

    # Fix tag text shapes: make smaller, ellipse, white bold centered
    for shape, letter in tag_text_shapes:
        sp = shape._element
        spPr = sp.find('p:spPr', NS)
        xfrm = spPr.find('a:xfrm', NS)
        off = xfrm.find('a:off', NS)
        ext = xfrm.find('a:ext', NS)

        old_x = int(off.get('x'))
        old_y = int(off.get('y'))
        old_w = int(ext.get('cx'))
        old_h = int(ext.get('cy'))

        center_x = old_x + old_w // 2
        center_y = old_y + old_h // 2
        new_x = center_x - TAG_SIZE // 2
        new_y = center_y - TAG_SIZE // 2

        off.set('x', str(new_x))
        off.set('y', str(new_y))
        ext.set('cx', str(TAG_SIZE))
        ext.set('cy', str(TAG_SIZE))

        # No fill on text shape (bg shape handles fill)
        existing_fill = spPr.find('a:solidFill', NS)
        if existing_fill is not None:
            spPr.remove(existing_fill)
        nofill = spPr.find('a:noFill', NS)
        if nofill is None:
            etree.SubElement(spPr, qn('a:noFill'))

        # Ellipse geometry
        geom = spPr.find('a:prstGeom', NS)
        if geom is not None:
            geom.set('prst', 'ellipse')

        # Remove border
        ln = spPr.find('a:ln', NS)
        if ln is not None:
            spPr.remove(ln)
        etree.SubElement(spPr, qn('a:ln'))

        # White bold text, centered
        for rpr in sp.findall('.//a:rPr', NS):
            rpr.set('b', '1')
            rpr.set('sz', '800')  # 8pt font for smaller circles
            for sf in rpr.findall('a:solidFill', NS):
                rpr.remove(sf)
            sf = etree.SubElement(rpr, qn('a:solidFill'))
            clr = etree.SubElement(sf, qn('a:srgbClr'))
            clr.set('val', 'FFFFFF')

        for bodyPr in sp.findall('.//a:bodyPr', NS):
            bodyPr.set('anchor', 'ctr')
            bodyPr.set('lIns', '0')
            bodyPr.set('tIns', '0')
            bodyPr.set('rIns', '0')
            bodyPr.set('bIns', '0')

        for pPr in sp.findall('.//a:pPr', NS):
            pPr.set('algn', 'ctr')

    # Style description labels (A., B., etc.)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if len(text) == 2 and text[0] in 'ABCDEFG' and text[1] == '.':
            for rpr in shape._element.findall('.//a:rPr', NS):
                rpr.set('b', '1')
                rpr.set('sz', '1100')
                for sf in rpr.findall('a:solidFill', NS):
                    rpr.remove(sf)
                sf = etree.SubElement(rpr, qn('a:solidFill'))
                clr = etree.SubElement(sf, qn('a:srgbClr'))
                clr.set('val', '3A9D6E')

    tag_count = len(tag_bg_shapes) + len(tag_text_shapes)
    if removed > 0 or tag_count > 0:
        print(f"  Slide {slide_idx+1}: removed {removed} connectors, fixed {len(tag_text_shapes)} tags")

# ============================================================
# STEP 3: Reorder slides & remove slide 17
# ============================================================
print("\n=== STEP 3: Reorder slides ===")

# Current order (1-indexed): 1,2,3,4,5,6,7,8,9,10,11,12,13,14,15,16,17,18
# Target order:  1,2,3,4, 5,6,7,8,16, 10,11,12,13,14, 9,15,18
# (remove 17)
# 0-indexed:     0,1,2,3, 4,5,6,7,15, 9,10,11,12,13, 8,14,17

new_order = [0, 1, 2, 3, 4, 5, 6, 7, 15, 9, 10, 11, 12, 13, 8, 14, 17]

# Get presentation.xml sldIdLst
pres_el = prs._element
sldIdLst = pres_el.find(qn('p:sldIdLst'))
sldIds = list(sldIdLst)

print(f"  Current slides: {len(sldIds)}")
print(f"  Target order (0-indexed): {new_order}")

# Reorder
new_sldIds = [sldIds[i] for i in new_order]

# Clear and re-add
for child in list(sldIdLst):
    sldIdLst.remove(child)
for sldId in new_sldIds:
    sldIdLst.append(sldId)

print(f"  New slide count: {len(new_sldIds)} (removed slide 17)")

# ============================================================
# STEP 4: Update page numbers on slides
# ============================================================
print("\n=== STEP 4: Update page numbers ===")

# After reorder, update the page number text on each slide
slides_after = list(prs.slides)
for idx, slide in enumerate(slides_after):
    page_num = idx + 1
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            # Page number is typically a standalone number at bottom-right
            if text.isdigit() and shape.left > 8000000 and shape.top > 4500000:
                # Update page number
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = str(page_num)
                print(f"  Slide {page_num}: updated page number from '{text}' to '{page_num}'")

# ============================================================
# SAVE
# ============================================================
output_path = "김윤찬_TRY-CATCH_화면_포트폴리오_v7_fixed.pptx"
prs.save(output_path)
print(f"\nSaved to: {output_path}")
print("\nNew slide order:")
print("  1. 표지")
print("  2. 기획 배경 및 의도")
print("  3. 서비스 기능 및 화면 리스트")
print("  4. 메뉴트리")
print("  5. Q&A 목록 (비로그인)")
print("  6. Q&A 작성")
print("  7. Q&A 상세")
print("  8. Q&A 수정")
print("  9. Q&A 작성/수정 세부 인터랙션")
print("  10. 마이페이지 프로필/홈")
print("  11. 개인정보 수정")
print("  12. 알림 목록")
print("  13. 활동 내역")
print("  14. 회원 탈퇴")
print("  15. 공통 푸터")
print("  16. 공통 드롭다운/상태")
print("  17. 공통 헤더 상태 변화")
