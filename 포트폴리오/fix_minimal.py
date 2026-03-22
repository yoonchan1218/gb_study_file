"""
최소 수정 - 기존 도형만 변경, 새 도형 추가 없음
"""
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

prs = Presentation("김윤찬_TRY-CATCH_화면_포트폴리오_v6.pptx")

# 1. Slide 3: 초록 오버레이 투명도
for clr in prs.slides[2]._element.iter(qn('a:srgbClr')):
    if clr.get('val') == 'DBEEE0':
        for old in clr.findall(qn('a:alpha')):
            clr.remove(old)
        a = etree.SubElement(clr, qn('a:alpha'))
        a.set('val', '45000')
print("1. Slide 3 done")

# 2. Slides 5-14: connector 제거 + 태그 수정
TAG = 146304

for si in range(4, 15):
    slide = prs.slides[si]
    spTree = slide._element.find('.//' + qn('p:spTree'))

    for c in spTree.findall(qn('p:cxnSp')):
        spTree.remove(c)

    img_right = 0
    for s in slide.shapes:
        try:
            _ = s.image
            if s.width > 2000000:
                img_right = s.left + s.width
        except:
            pass
    if img_right == 0:
        continue

    nx = img_right - TAG - 80000

    for s in slide.shapes:
        sp = s._element
        spPr = sp.find(qn('p:spPr'))
        if spPr is None:
            continue

        hit = False
        if not s.has_text_frame:
            f = spPr.find(qn('a:solidFill'))
            if f is not None:
                c2 = f.find(qn('a:srgbClr'))
                if c2 is not None and c2.get('val') == '3A9D6E' and s.width == 256032:
                    hit = True
        elif s.text_frame.text.strip() in list('ABCDEFG') and s.width == 256032:
            hit = True
            txt = s.text_frame.text.strip()
            for rpr in sp.iter(qn('a:rPr')):
                rpr.set('sz', '700')
                rpr.set('b', '1')
                for sf in rpr.findall(qn('a:solidFill')):
                    rpr.remove(sf)
                sf = etree.SubElement(rpr, qn('a:solidFill'))
                cc = etree.SubElement(sf, qn('a:srgbClr'))
                cc.set('val', 'FFFFFF')
            for bp in sp.iter(qn('a:bodyPr')):
                bp.set('anchor', 'ctr')
                bp.set('lIns', '0')
                bp.set('tIns', '0')
                bp.set('rIns', '0')
                bp.set('bIns', '0')
            for pp in sp.iter(qn('a:pPr')):
                pp.set('algn', 'ctr')

        if hit:
            cy = s.top + s.height // 2
            s.left = nx
            s.top = cy - TAG // 2
            s.width = TAG
            s.height = TAG
            g = spPr.find(qn('a:prstGeom'))
            if g is not None:
                g.set('prst', 'ellipse')
            ln = spPr.find(qn('a:ln'))
            if ln is not None:
                spPr.remove(ln)
            etree.SubElement(spPr, qn('a:ln'))

    print(f"  Slide {si+1} done")

out = "김윤찬_TRY-CATCH_화면_포트폴리오_v7_clean.pptx"
prs.save(out)
import os
print(f"\nSaved: {out} ({os.path.getsize(out):,} bytes)")
