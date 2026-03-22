"""
안전한 수정만 적용 - 슬라이드 삭제/재정렬 없음
1. Slide 3: 초록 오버레이 투명도
2. Slides 5-14: connector 선 제거, 태그 작은 원형 + 이미지 안쪽 배치
3. Slides 16,18: 헤더 추가
"""
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'

prs = Presentation("김윤찬_TRY-CATCH_화면_포트폴리오_v6.pptx")

# ── 1. Slide 3: 초록 오버레이 투명도 ──
print("1. Slide 3 투명도")
for clr in prs.slides[2]._element.iter(qn('a:srgbClr')):
    if clr.get('val') == 'DBEEE0':
        for old in clr.findall(qn('a:alpha')):
            clr.remove(old)
        a = etree.SubElement(clr, qn('a:alpha'))
        a.set('val', '45000')

# ── 2. Slides 5-14: connector 제거 + 태그 수정 ──
print("2. Connector 제거 + 태그 수정")
TAG_SIZE = 146304

for si in range(4, 15):
    slide = prs.slides[si]
    spTree = slide._element.find('.//' + qn('p:spTree'))

    # connector 제거
    for c in spTree.findall(qn('p:cxnSp')):
        spTree.remove(c)

    # 메인 이미지 우측 경계 찾기
    img_right = 0
    for s in slide.shapes:
        try:
            _ = s.image
            if s.width > 2000000:
                img_right = s.left + s.width
        except:
            pass

    if img_right == 0:
        continue

    new_tag_x = img_right - TAG_SIZE - 80000

    for s in slide.shapes:
        sp = s._element
        spPr = sp.find(qn('p:spPr'))
        if spPr is None:
            continue

        is_tag_bg = False
        is_tag_text = False

        if not s.has_text_frame:
            fill = spPr.find(qn('a:solidFill'))
            if fill is not None:
                clr = fill.find(qn('a:srgbClr'))
                if clr is not None and clr.get('val') == '3A9D6E' and s.width == 256032:
                    is_tag_bg = True
        else:
            txt = s.text_frame.text.strip()
            if len(txt) == 1 and txt in 'ABCDEFG' and s.width == 256032:
                is_tag_text = True

        if is_tag_bg or is_tag_text:
            cy = s.top + s.height // 2
            s.left = new_tag_x
            s.top = cy - TAG_SIZE // 2
            s.width = TAG_SIZE
            s.height = TAG_SIZE
            geom = spPr.find(qn('a:prstGeom'))
            if geom is not None:
                geom.set('prst', 'ellipse')
            ln = spPr.find(qn('a:ln'))
            if ln is not None:
                spPr.remove(ln)
            etree.SubElement(spPr, qn('a:ln'))

        if is_tag_text:
            for rpr in sp.iter(qn('a:rPr')):
                rpr.set('sz', '700')
                rpr.set('b', '1')
                for sf in rpr.findall(qn('a:solidFill')):
                    rpr.remove(sf)
                sf = etree.SubElement(rpr, qn('a:solidFill'))
                c2 = etree.SubElement(sf, qn('a:srgbClr'))
                c2.set('val', 'FFFFFF')
            for bp in sp.iter(qn('a:bodyPr')):
                bp.set('anchor', 'ctr')
                bp.set('lIns', '0')
                bp.set('tIns', '0')
                bp.set('rIns', '0')
                bp.set('bIns', '0')
            for pp in sp.iter(qn('a:pPr')):
                pp.set('algn', 'ctr')

        # A./B. 라벨 스타일
        if s.has_text_frame:
            txt = s.text_frame.text.strip()
            if len(txt) == 2 and txt[0] in 'ABCDEFG' and txt[1] == '.':
                for rpr in sp.iter(qn('a:rPr')):
                    rpr.set('b', '1')
                    for sf in rpr.findall(qn('a:solidFill')):
                        rpr.remove(sf)
                    sf = etree.SubElement(rpr, qn('a:solidFill'))
                    c2 = etree.SubElement(sf, qn('a:srgbClr'))
                    c2.set('val', '3A9D6E')

    print(f"  Slide {si+1} done")

# ── 3. Slides 16, 18: 헤더 추가 ──
print("3. 이미지 슬라이드 헤더 추가")

def add_header(slide, subtitle):
    spTree = slide._element.find('.//' + qn('p:spTree'))
    # 이미지 축소
    for s in slide.shapes:
        try:
            _ = s.image
            s.top = 457200
            s.height = 4686300
        except:
            pass

    def make_rect(parent, sid, name, x, y, w, h, fill_val=None):
        sp = etree.SubElement(parent, qn('p:sp'))
        nv = etree.SubElement(sp, qn('p:nvSpPr'))
        cp = etree.SubElement(nv, qn('p:cNvPr'))
        cp.set('id', str(sid))
        cp.set('name', name)
        etree.SubElement(nv, qn('p:cNvSpPr'))
        etree.SubElement(nv, qn('p:nvPr'))
        pr = etree.SubElement(sp, qn('p:spPr'))
        xf = etree.SubElement(pr, qn('a:xfrm'))
        o = etree.SubElement(xf, qn('a:off'))
        o.set('x', str(x)); o.set('y', str(y))
        e = etree.SubElement(xf, qn('a:ext'))
        e.set('cx', str(w)); e.set('cy', str(h))
        g = etree.SubElement(pr, qn('a:prstGeom'))
        g.set('prst', 'rect')
        etree.SubElement(g, qn('a:avLst'))
        if fill_val:
            f = etree.SubElement(pr, qn('a:solidFill'))
            c = etree.SubElement(f, qn('a:srgbClr'))
            c.set('val', fill_val)
        else:
            etree.SubElement(pr, qn('a:noFill'))
        etree.SubElement(pr, qn('a:ln'))
        return sp

    def add_text(sp, text, sz, bold, color):
        tb = etree.SubElement(sp, qn('p:txBody'))
        bp = etree.SubElement(tb, qn('a:bodyPr'))
        bp.set('wrap', 'square'); bp.set('lIns', '0'); bp.set('tIns', '0')
        bp.set('rIns', '0'); bp.set('bIns', '0'); bp.set('anchor', 'ctr')
        etree.SubElement(tb, qn('a:lstStyle'))
        p = etree.SubElement(tb, qn('a:p'))
        pp = etree.SubElement(p, qn('a:pPr'))
        pp.set('indent', '0'); pp.set('marL', '0')
        etree.SubElement(pp, qn('a:buNone'))
        r = etree.SubElement(p, qn('a:r'))
        rp = etree.SubElement(r, qn('a:rPr'))
        rp.set('lang', 'ko-KR'); rp.set('sz', str(sz)); rp.set('dirty', '0')
        if bold:
            rp.set('b', '1')
        sf = etree.SubElement(rp, qn('a:solidFill'))
        sc = etree.SubElement(sf, qn('a:srgbClr'))
        sc.set('val', color)
        for face in [qn('a:latin'), qn('a:ea')]:
            fn = etree.SubElement(rp, face)
            fn.set('typeface', '\ub9d1\uc740 \uace0\ub515')
        t = etree.SubElement(r, qn('a:t'))
        t.text = text

    # 좌측 어센트 바
    bar = make_rect(spTree, 900, 'Bar', 0, 0, 54864, 457200, '333333')
    etree.SubElement(bar, qn('p:txBody')).append(etree.Element(qn('a:bodyPr')))

    # 타이틀
    title = make_rect(spTree, 901, 'Title', 182880, 45720, 8229600, 320040)
    add_text(title, 'TRY-CATCH    \ud654\uba74 \ub0b4 \uc601\uc5ed \uc124\uba85', 1400, True, '333333')

    # 서브타이틀
    sub = make_rect(spTree, 902, 'Sub', 182880, 384048, 4572000, 73152)
    add_text(sub, subtitle, 1000, False, '666666')

for idx, sub_text in [(15, '\u300cQ&A \uc791\uc131/\uc218\uc815 \uc138\ubd80 \uc778\ud130\ub799\uc158\u300d'),
                       (17, '\u300c\uacf5\ud1b5 \uc0c1\ud0dc \ubcc0\ud654: \ub85c\uadf8\uc778/\uc54c\ub9bc/\uac80\uc0c9\u300d')]:
    add_header(prs.slides[idx], sub_text)
    print(f"  Slide {idx+1} header added")

# ── 저장 ──
out = "김윤찬_TRY-CATCH_화면_포트폴리오_v7_final.pptx"
prs.save(out)
import os
print(f"\nSaved: {out} ({os.path.getsize(out):,} bytes)")
print("슬라이드 17 삭제 및 순서 변경은 PowerPoint에서 직접 해주세요.")
print("  - 슬라이드 17(체험 지원) 우클릭 → 삭제")
print("  - 슬라이드 16(Q&A 세부)을 8번 뒤로 드래그")
print("  - 슬라이드 9(푸터)를 14번 뒤로 드래그")
print("  - 슬라이드 15(드롭다운)를 푸터 뒤로 드래그")
print("  - 슬라이드 18(공통 헤더)을 드롭다운 뒤로 드래그")
