# -*- coding: utf-8 -*-
"""
Rebuild slides 9-18 of the portfolio PPTX to match slides 5-8 pattern.
"""
import os
import sys
import copy
from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image

# ── Paths ──
BASE = r'C:\Users\pigch\Desktop\gb_0090_kyc\포트폴리오'
SRC = os.path.join(BASE, '김윤찬_TRY-CATCH_화면_포트폴리오_v6.pptx')
DST = os.path.join(BASE, '김윤찬_TRY-CATCH_화면_포트폴리오_FINAL.pptx')
SHOTS = os.path.join(BASE, 'screenshots')

# ── Constants from template analysis ──
SLIDE_W = 9144000
SLIDE_H = 5143500
FONT_NAME = '맑은 고딕'

# Header elements
ACCENT_BAR = dict(left=0, top=0, width=54864, height=457200, fill_color='333333')
TITLE = dict(left=182880, top=45720, width=4572000, height=320040,
             text='TRY-CATCH    화면 내 영역 설명', size=Pt(14), bold=True, color='333333')
PAGE_NUM = dict(left=8686800, top=4846320, width=274320, height=182880, size=Pt(8), color='888888')

# Description area constants
DESC_HEADER_H = 228600
SEP_H = 9144
TAG_LABEL_W = 228600
TAG_LABEL_H = 182880
TAG_TITLE_H = 137160
TAG_DETAIL_H = 338328
TAG_SPACING = 502920   # Y increment between tags

# Tag circle
CIRCLE_SIZE = 256032
CIRCLE_FILL = '3A9D6E'

# Image area
IMG_TOP = 658368
IMG_AREA_TOP = 640080
IMG_LEFT = 137160
BG_LEFT = 118872
BG_TOP = 640080

TAG_LETTERS = 'ABCDEFGHIJKLMNOP'


def get_image_dims(path):
    """Return (w, h) in pixels."""
    with Image.open(path) as im:
        return im.size


def fit_image(img_w, img_h, max_w, max_h):
    """Return (w, h) in EMU preserving aspect ratio."""
    ratio = min(max_w / img_w, max_h / img_h)
    return int(img_w * ratio), int(img_h * ratio)


def clear_slide(slide):
    """Remove ALL shapes from a slide's spTree, keeping group envelope."""
    spTree = slide.shapes._spTree
    nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
             'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    # Keep nvGrpSpPr and grpSpPr, remove everything else
    to_remove = []
    for child in spTree:
        tag = etree.QName(child).localname
        if tag not in ('nvGrpSpPr', 'grpSpPr'):
            to_remove.append(child)
    for child in to_remove:
        spTree.remove(child)


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=None):
    """Add a rectangle shape with solid fill."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()  # no line by default
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(fill_color)
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = RGBColor.from_string(border_color)
        shape.line.width = border_width or Pt(0.5)
    return shape


def add_textbox(slide, left, top, width, height, text, font_size, bold=None,
                color='333333', font_name=FONT_NAME, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """Add a text box with specified formatting."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.background()
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    if bold is not None:
        run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)
    # Set vertical alignment
    txBody = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
    if txBody is not None:
        if anchor == MSO_ANCHOR.MIDDLE:
            txBody.set('anchor', 'ctr')
        elif anchor == MSO_ANCHOR.BOTTOM:
            txBody.set('anchor', 'b')
    return shape


def add_circle(slide, left, top, size, fill_color, text, font_size, text_color='FFFFFF'):
    """Add a filled circle (bg) + text circle pair like the template."""
    # Background circle (solid fill)
    bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor.from_string(fill_color)
    bg.line.fill.background()

    # Text circle (transparent, with centered text)
    txt = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    txt.fill.background()
    txt.line.fill.background()
    tf = txt.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = font_size
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string(text_color)
    # Center vertically
    txBody = txt._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
    if txBody is not None:
        txBody.set('anchor', 'ctr')
    return bg, txt


def add_separator(slide, left, top, width):
    """Add a thin separator line (as a rectangle)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, SEP_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string('CCCCCC')
    shape.line.fill.background()
    return shape


def add_connector_line(slide, x1, y1, x2, y2):
    """Add a connector line from circle to description area."""
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # MSO_CONNECTOR_TYPE.STRAIGHT = 1
    connector.line.color.rgb = RGBColor.from_string('3A9D6E')
    connector.line.width = Pt(0.75)
    connector.line.dash_style = 2  # dash
    return connector


def build_standard_slide(slide, subtitle_text, image_path, desc_header_text,
                         tags, page_num, tag_y_fractions=None):
    """
    Build a standard single-image slide matching the 5-8 pattern.

    tags: list of (label_letter, title, detail)
    tag_y_fractions: list of floats [0..1] indicating where on the image each tag circle sits
    """
    clear_slide(slide)

    # -- Accent bar --
    add_rect(slide, **ACCENT_BAR)

    # -- Title --
    add_textbox(slide, TITLE['left'], TITLE['top'], TITLE['width'], TITLE['height'],
                TITLE['text'], TITLE['size'], bold=True, color=TITLE['color'])

    # -- Subtitle --
    add_textbox(slide, 182880, 384048, 4572000, 201168,
                f'「{subtitle_text}」', Pt(11), bold=True, color='333333')

    # -- Load image, compute sizing --
    pw, ph = get_image_dims(image_path)

    # Determine available image area height
    img_area_h = SLIDE_H - IMG_TOP - 100000  # leave ~1cm bottom margin
    # max image height = 4114800 (from template)
    img_area_h = min(img_area_h, 4114800)

    # Number of tags determines description width needed
    n_tags = len(tags)

    # Determine image width: for fewer tags, image can be wider
    # Template pattern: desc area starts at image_right + gap
    # We want desc area width ~ 4000000-5200000 EMU depending on text
    desc_area_w = 4100000
    gap = 100000
    max_img_w = SLIDE_W - desc_area_w - gap - IMG_LEFT - 200000

    # Fit image
    iw, ih = fit_image(pw, ph, max_img_w, img_area_h)

    # Image background box
    bg_w = iw + 36576  # small padding
    bg_h = ih + 36576
    add_rect(slide, BG_LEFT, BG_TOP, bg_w, bg_h, fill_color='FAFAFA', border_color='DDDDDD', border_width=Pt(0.5))

    # Image
    pic = slide.shapes.add_picture(image_path, IMG_LEFT, IMG_TOP, iw, ih)

    # Image right edge
    img_right = IMG_LEFT + iw

    # Description area x
    desc_x = img_right + 200000
    desc_content_x = desc_x + TAG_LABEL_W
    desc_content_w = SLIDE_W - desc_content_x - 91440  # right margin

    # -- Description header --
    add_textbox(slide, desc_x, 384048, desc_content_w + TAG_LABEL_W, DESC_HEADER_H,
                f'{desc_header_text}', Pt(10), bold=True, color='006E3F')

    # -- Separator --
    add_separator(slide, desc_x, 594360, desc_content_w + TAG_LABEL_W)

    # -- Tags --
    # Tag Y positions on description side (start at 685800, increment TAG_SPACING)
    tag_start_y = 685800

    # Default tag_y_fractions if not specified
    if tag_y_fractions is None:
        tag_y_fractions = [i / max(n_tags - 1, 1) for i in range(n_tags)]

    for i, (letter, title, detail) in enumerate(tags):
        desc_y = tag_start_y + i * TAG_SPACING

        # Tag label "A."
        add_textbox(slide, desc_x, desc_y, TAG_LABEL_W, TAG_LABEL_H,
                    f'{letter}.', Pt(10), bold=True, color='3A9D6E')

        # Tag title
        add_textbox(slide, desc_content_x, desc_y, desc_content_w, TAG_TITLE_H,
                    title, Pt(9), bold=True, color='333333')

        # Tag detail
        add_textbox(slide, desc_content_x, desc_y + TAG_TITLE_H, desc_content_w, TAG_DETAIL_H,
                    detail, Pt(7), bold=False, color='666666')

        # Circle on image (near right edge)
        circle_x = img_right - CIRCLE_SIZE - 18288  # inside right edge
        # Circle Y on image based on fraction
        frac = tag_y_fractions[i] if i < len(tag_y_fractions) else 0.5
        circle_y = int(IMG_TOP + frac * (ih - CIRCLE_SIZE))
        circle_y = max(circle_y, IMG_TOP)
        circle_y = min(circle_y, IMG_TOP + ih - CIRCLE_SIZE)

        add_circle(slide, circle_x, circle_y, CIRCLE_SIZE, CIRCLE_FILL,
                   letter, Pt(10))

        # Connector line from circle to description
        line_start_x = circle_x + CIRCLE_SIZE
        line_start_y = circle_y + CIRCLE_SIZE // 2
        line_end_x = desc_x
        line_end_y = desc_y + TAG_LABEL_H // 2
        add_connector_line(slide, line_start_x, line_start_y, line_end_x, line_end_y)

    # -- Page number --
    add_textbox(slide, PAGE_NUM['left'], PAGE_NUM['top'], PAGE_NUM['width'], PAGE_NUM['height'],
                str(page_num), PAGE_NUM['size'], color=PAGE_NUM['color'], align=PP_ALIGN.RIGHT)


def build_multi_image_slide(slide, subtitle_text, image_paths, desc_header_text,
                            tags, page_num):
    """
    Build slide 15 - multi-image layout.
    Multiple screenshots arranged in a grid, description on the right bottom.
    """
    clear_slide(slide)

    # -- Accent bar --
    add_rect(slide, **ACCENT_BAR)

    # -- Title --
    add_textbox(slide, TITLE['left'], TITLE['top'], TITLE['width'], TITLE['height'],
                TITLE['text'], TITLE['size'], bold=True, color=TITLE['color'])

    # -- Subtitle --
    add_textbox(slide, 182880, 384048, 4572000, 201168,
                f'「{subtitle_text}」', Pt(11), bold=True, color='333333')

    # Layout: 2x2 grid of images on left side, descriptions on right
    img_area_left = 137160
    img_gap = 54864
    n_images = len(image_paths)

    # Arrange in 2 columns, 2 rows
    cols = 2
    rows = (n_images + 1) // 2

    img_area_top = 685800
    img_area_bottom = 4800000
    img_area_right = 5000000  # leave room for descriptions

    cell_w = (img_area_right - img_area_left - img_gap * (cols - 1)) // cols
    cell_h = (img_area_bottom - img_area_top - img_gap * (rows - 1)) // rows

    image_positions = []
    for idx, img_path in enumerate(image_paths):
        if not os.path.exists(img_path):
            continue
        row = idx // cols
        col = idx % cols

        cx = img_area_left + col * (cell_w + img_gap)
        cy = img_area_top + row * (cell_h + img_gap)

        pw, ph = get_image_dims(img_path)
        iw, ih = fit_image(pw, ph, cell_w, cell_h)

        # Background box
        add_rect(slide, cx - 9144, cy - 9144, iw + 18288, ih + 18288,
                 fill_color='FAFAFA', border_color='DDDDDD', border_width=Pt(0.5))

        pic = slide.shapes.add_picture(img_path, cx, cy, iw, ih)
        image_positions.append((cx, cy, iw, ih))

    # Description area on right side
    desc_x = 5200000
    desc_content_x = desc_x + TAG_LABEL_W
    desc_content_w = SLIDE_W - desc_content_x - 91440

    # Description header
    add_textbox(slide, desc_x, 384048, desc_content_w + TAG_LABEL_W, DESC_HEADER_H,
                f'{desc_header_text}', Pt(10), bold=True, color='006E3F')

    # Separator
    add_separator(slide, desc_x, 594360, desc_content_w + TAG_LABEL_W)

    # Tags
    tag_start_y = 685800
    for i, (letter, title, detail) in enumerate(tags):
        desc_y = tag_start_y + i * TAG_SPACING

        add_textbox(slide, desc_x, desc_y, TAG_LABEL_W, TAG_LABEL_H,
                    f'{letter}.', Pt(10), bold=True, color='3A9D6E')

        add_textbox(slide, desc_content_x, desc_y, desc_content_w, TAG_TITLE_H,
                    title, Pt(9), bold=True, color='333333')

        add_textbox(slide, desc_content_x, desc_y + TAG_TITLE_H, desc_content_w, TAG_DETAIL_H,
                    detail, Pt(7), bold=False, color='666666')

    # Page number
    add_textbox(slide, PAGE_NUM['left'], PAGE_NUM['top'], PAGE_NUM['width'], PAGE_NUM['height'],
                str(page_num), PAGE_NUM['size'], color=PAGE_NUM['color'], align=PP_ALIGN.RIGHT)


def build_fullimage_slide_with_header(slide, subtitle_text, page_num):
    """
    For slides 16, 18: keep the existing full image but add header elements on top.
    Resize image to start below header.
    """
    # Find the existing image
    img_shape = None
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            img_shape = shape
            break

    if img_shape is None:
        return

    # Extract image blob before clearing
    img_blob = img_shape.image.blob
    img_content_type = img_shape.image.content_type

    # Get image dimensions for aspect ratio
    import io
    with Image.open(io.BytesIO(img_blob)) as im:
        pw, ph = im.size

    clear_slide(slide)

    # -- Accent bar --
    add_rect(slide, **ACCENT_BAR)

    # -- Title --
    add_textbox(slide, TITLE['left'], TITLE['top'], TITLE['width'], TITLE['height'],
                TITLE['text'], TITLE['size'], bold=True, color=TITLE['color'])

    # -- Subtitle --
    add_textbox(slide, 182880, 384048, 4572000, 201168,
                f'「{subtitle_text}」', Pt(11), bold=True, color='333333')

    # -- Image below header --
    img_top = 500000
    img_available_h = SLIDE_H - img_top - 50000
    img_available_w = SLIDE_W - 100000

    iw, ih = fit_image(pw, ph, img_available_w, img_available_h)
    # Center horizontally
    img_left = (SLIDE_W - iw) // 2

    # Save blob to temp file and add
    import tempfile
    ext = '.png' if 'png' in img_content_type else '.jpg'
    tmp = tempfile.NamedTemporaryFile(suffix=ext, delete=False)
    tmp.write(img_blob)
    tmp.close()

    slide.shapes.add_picture(tmp.name, img_left, img_top, iw, ih)
    os.unlink(tmp.name)

    # -- Page number --
    add_textbox(slide, PAGE_NUM['left'], PAGE_NUM['top'], PAGE_NUM['width'], PAGE_NUM['height'],
                str(page_num), PAGE_NUM['size'], color=PAGE_NUM['color'], align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════
# Main
# ══════════════════════════════════════════════
def main():
    prs = Presentation(SRC)
    slides = prs.slides

    # ── Slide 9: Footer ──
    build_standard_slide(
        slides[8],
        subtitle_text='공통 푸터 (Footer)',
        image_path=os.path.join(SHOTS, 's9_footer.png'),
        desc_header_text='공통 푸터 (Footer) 화면 영역 설명',
        tags=[
            ('A', '제휴 광고', '체험 관련 외부 광고, 배너/채팅 안내'),
            ('B', '공지사항', '사이트 주요 공지사항 한줄 표시'),
            ('C', '사이트 정책 링크', '회사소개, 이용약관, 개인정보처리방침, 고객센터'),
            ('D', '연락처/저작권', '전화, 운영시간, FAX, 이메일, ©표시'),
        ],
        page_num=9,
        tag_y_fractions=[0.60, 0.72, 0.82, 0.92],  # footer is at bottom
    )

    # ── Slide 10: 마이페이지 홈/프로필 ──
    build_standard_slide(
        slides[9],
        subtitle_text='마이페이지 홈/프로필',
        image_path=os.path.join(SHOTS, 's10_mypage_home.png'),
        desc_header_text='마이페이지 홈/프로필 화면 영역 설명',
        tags=[
            ('A', '사이드 네비게이션', '개인회원 홈, 기술블로그, 체험지원관리, QnA, 회원정보 관리'),
            ('B', '프로필 카드', '프로필 사진 업로드, 닉네임, 레벨, 포인트, 별 충전'),
            ('C', '최근 체험 슬라이더', '추천 체험 프로그램 Swiper 슬라이드 (좌우 넘김)'),
            ('D', '인기글/최근글', '인기 Q&A 목록, 태그/카테고리 표시'),
            ('E', '최근 체험공고/TOP100', '최근 프로그램, TOP100 추천 카드'),
            ('F', '포인트 매칭 (선택)', '최근 별 내역/스크랩 게시물 전환'),
        ],
        page_num=10,
        tag_y_fractions=[0.05, 0.15, 0.35, 0.52, 0.72, 0.88],
    )

    # ── Slide 11: 개인정보 수정 ──
    build_standard_slide(
        slides[10],
        subtitle_text='개인정보 수정',
        image_path=os.path.join(SHOTS, 's11_edit_info.png'),
        desc_header_text='개인정보 수정 화면 영역 설명',
        tags=[
            ('A', '수정화면 타이틀', "'회원정보 수정' 타이틀"),
            ('B', '아이디 (읽기전용)', '회원 아이디 표시 (수정 불가)'),
            ('C', '기본 정보 수정', '이름, 생년월일 드롭다운, 성별 선택, 학력'),
            ('D', '연락처 수정', '휴대폰, 일반전화, 이메일(@도메인)'),
            ('E', '수정하기/취소 버튼', '폼 전송, 유효성 검증'),
        ],
        page_num=11,
        tag_y_fractions=[0.05, 0.18, 0.40, 0.65, 0.90],
    )

    # ── Slide 12: 알림 목록 ──
    build_standard_slide(
        slides[11],
        subtitle_text='알림 목록',
        image_path=os.path.join(SHOTS, 's12_notification.png'),
        desc_header_text='알림 목록 화면 영역 설명',
        tags=[
            ('A', '사이드 네비게이션', '개인회원 홈, 기술블로그, 체험지원관리, QnA, 회원정보 관리'),
            ('B', '수신 알림 타이틀', "'수신 알림', 최근 30일 알림 안내"),
            ('C', '카테고리 필터', '전체/체험지원/체험안내/기술블로그 드롭다운 필터링'),
            ('D', '알림 목록/빈 상태', "날짜별 그룹핑 또는 '알림 없음' 안내"),
        ],
        page_num=12,
        tag_y_fractions=[0.05, 0.15, 0.22, 0.45],
    )

    # ── Slide 13: 활동 내역 (체험지원 현황) ──
    build_standard_slide(
        slides[12],
        subtitle_text='활동 내역 (체험지원 현황)',
        image_path=os.path.join(SHOTS, 's13_experience.png'),
        desc_header_text='활동 내역 (체험지원 현황) 화면 영역 설명',
        tags=[
            ('A', '상태 요약 카드', '참여완료/참여중/심사완료/심사중 카운트'),
            ('B', '날짜/상태 검색', '조회기간, 날짜, 진행/합격/불합격여부 필터'),
            ('C', '지원 현황 목록', '기업명, 프로그램, 지원분야별, 상태표시 버튼'),
            ('D', '페이지네이션', '지원 목록 페이지 이동'),
            ('E', '유의사항/푸터', '유의사항 5가지, 연락 정보, 회사 정보'),
        ],
        page_num=13,
        tag_y_fractions=[0.08, 0.22, 0.45, 0.70, 0.88],
    )

    # ── Slide 14: 회원 탈퇴 ──
    build_standard_slide(
        slides[13],
        subtitle_text='회원 탈퇴',
        image_path=os.path.join(SHOTS, 's14_unsubscribe.png'),
        desc_header_text='회원 탈퇴 화면 영역 설명',
        tags=[
            ('A', '탈퇴 화면 진입', "로고, '회원정보 관리', 수정/탈퇴 탭"),
            ('B', '유의사항 목록', '계정 복구 불가, ID 재사용 불가, 동시 탈퇴, 데이터 삭제'),
            ('C', '본인 확인 폼', '아이디(읽기전용), 탈퇴요청자 이름 입력'),
            ('D', '동의 및 탈퇴 버튼', '체크박스 선택 후 탈퇴하기 활성화'),
        ],
        page_num=14,
        tag_y_fractions=[0.05, 0.30, 0.60, 0.85],
    )

    # ── Slide 15: Multi-image (Header states) ──
    build_multi_image_slide(
        slides[14],
        subtitle_text='공통 Header 상태 + 드롭다운',
        image_paths=[
            os.path.join(SHOTS, 's15_loggedout.png'),
            os.path.join(SHOTS, 's15_profile_dropdown.png'),
            os.path.join(SHOTS, 's15_search_autocomplete.png'),
            os.path.join(SHOTS, 's15_alarm_dropdown.png'),
        ],
        desc_header_text='공통 Header 상태 + 드롭다운 화면 영역 설명',
        tags=[
            ('A', 'Header 드롭다운', '서비스소개/체험정보/기술블로그/커뮤니티, hover 시 하위 메뉴 표시'),
            ('B', '비회원서비스 드롭다운', '비로그인 시 비회원서비스 버튼 클릭 시 체험 신청하기/개인회원 홈 안내'),
            ('C', '프로필 드롭다운', '로그인 후 회원 이름 클릭 시 개인회원홈/기술블로그/체험지원/개인정보/로그아웃'),
            ('D', '검색 자동완성', '키워드 자동완성 드롭다운, 자동완성 끄기/닫기 옵션 지원'),
            ('E', '알림 드롭다운', '새 알림 표시, 클릭 시 알림 목록'),
        ],
        page_num=15,
    )

    # ── Slide 16: Full-image with header ──
    build_fullimage_slide_with_header(
        slides[15],
        subtitle_text='Q&A 작성/수정 세부 인터랙션',
        page_num=16,
    )

    # ── Slide 17: keep as-is but add header ──
    # (Slide 17 is not mentioned in steps, but we should handle it consistently)
    # The instructions say slides 9-18 must match. Slide 17 is a full-image slide too.
    # Let's add header to it as well.
    build_fullimage_slide_with_header(
        slides[16],
        subtitle_text='Q&A 세부 인터랙션 (계속)',
        page_num=17,
    )

    # ── Slide 18: Full-image with header ──
    build_fullimage_slide_with_header(
        slides[17],
        subtitle_text='공통 상태 변화: 로그인 / 알림 / 검색',
        page_num=18,
    )

    # Save
    prs.save(DST)
    print(f'Saved to {DST}')


if __name__ == '__main__':
    main()
