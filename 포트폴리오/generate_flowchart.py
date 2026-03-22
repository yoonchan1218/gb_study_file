"""
TryCatch Backend Flowchart Generator v4
- DB icon: template image only (no text overlay)
- REST: positioned exactly at Controller level on divider
- Transaction: on ALL slides
- Flow name labels removed (page title is enough)
- Table names cleanly above DB icon, no overlap
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn
from lxml import etree
import os

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
DB_ICON_PATH = os.path.join(SCRIPT_DIR, "db_group_icon.png")

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
GRAY_BG = RGBColor(217, 217, 217)

SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)

C_LEFT = Inches(1.3)
C_TOP = Inches(1.05)
C_WIDTH = Inches(10.7)
C_HEIGHT = Inches(5.85)
HEADER_H = Inches(0.38)
DIVIDER_X = C_LEFT + Inches(6.0)

L_CENTER_X = C_LEFT + Inches(3.2)
L_BOX_W = Inches(3.2)

R_CENTER_X = DIVIDER_X + Inches(2.2)
R_BOX_W = Inches(2.0)

BOX_H = Inches(0.268)
ARROW_GAP = Inches(0.204)
CONNECT_R = Inches(0.13)


def _set_text(shape, text, font_size=9, bold=False, color=BLACK, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    tf.margin_left = Pt(3)
    tf.margin_right = Pt(3)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "맑은 고딕"


def add_textbox(slide, left, top, width, height, text,
                font_size=10, bold=False, align=PP_ALIGN.CENTER, color=BLACK):
    tb = slide.shapes.add_textbox(left, top, width, height)
    _set_text(tb, text, font_size, bold, color, align)
    return tb


def add_process_box(slide, left, top, width, height, text, font_size=9):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.FLOWCHART_PROCESS, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BLACK
    shape.line.width = Pt(0.75)
    _set_text(shape, text, font_size)
    return shape


def add_document_box(slide, left, top, width, height, text, font_size=9):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.FLOWCHART_DOCUMENT, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BLACK
    shape.line.width = Pt(0.75)
    _set_text(shape, text, font_size)
    return shape


def add_decision_diamond(slide, left, top, width, height, text, font_size=8):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.FLOWCHART_DECISION, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BLACK
    shape.line.width = Pt(0.75)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.margin_left = Pt(14)
    tf.margin_right = Pt(14)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = BLACK
        run.font.name = "맑은 고딕"
    return shape


def add_ellipse_btn(slide, left, top, width, height, text, font_size=8):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BLACK
    shape.line.width = Pt(0.75)
    _set_text(shape, text, font_size)
    return shape


def add_connect_circle(slide, cx, cy):
    r = CONNECT_R
    shape = slide.shapes.add_shape(
        MSO_SHAPE.FLOWCHART_CONNECTOR, cx - r, cy - r, r * 2, r * 2)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLACK
    shape.line.color.rgb = BLACK
    shape.line.width = Pt(0.5)
    return shape


def add_db_icon(slide, left, top):
    pic = slide.shapes.add_picture(
        DB_ICON_PATH, left, top, Inches(0.50), Inches(0.50))
    return pic


def _add_tail_arrow(conn):
    cxn_elem = conn._element
    spPr = cxn_elem.find(qn('p:spPr'))
    if spPr is None:
        return
    ln = spPr.find(qn('a:ln'))
    if ln is None:
        ln = etree.SubElement(spPr, qn('a:ln'))
    for te in ln.findall(qn('a:tailEnd')):
        ln.remove(te)
    tailEnd = etree.SubElement(ln, qn('a:tailEnd'))
    tailEnd.set('type', 'triangle')
    tailEnd.set('w', 'med')
    tailEnd.set('len', 'med')


def add_arrow_down(slide, x, y_from, y_to):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, x, y_from, x, y_to)
    conn.line.color.rgb = BLACK
    conn.line.width = Pt(0.75)
    _add_tail_arrow(conn)
    return conn


def add_line(slide, x1, y1, x2, y2, dashed=False):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = BLACK
    conn.line.width = Pt(0.75)
    if dashed:
        conn.line.dash_style = 4
    return conn


def add_line_arrow(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = BLACK
    conn.line.width = Pt(0.75)
    _add_tail_arrow(conn)
    return conn


def calc_box_h(items):
    n = len(items) if isinstance(items, list) else 1
    if n <= 1:
        return BOX_H
    return Inches(0.12 + 0.155 * n)


def create_slide(prs, page_num, category, page_title, flow_name,
                 db_tables, html_file,
                 mapper_xmls, mapper_javas, daos,
                 service, controller,
                 has_button=False, button_text="",
                 frontend_html=None,
                 has_decision=False, decision_text="",
                 exception_handler=None,
                 rest_label="REST"):

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if frontend_html is None:
        frontend_html = html_file

    # ── TITLE ──
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.28), Inches(0.06), Inches(0.48))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLACK
    bar.line.fill.background()

    add_textbox(slide, Inches(0.68), Inches(0.28), Inches(5), Inches(0.48),
                category, font_size=22, align=PP_ALIGN.LEFT)

    add_textbox(slide, Inches(12.5), Inches(0.12), Inches(0.6), Inches(0.3),
                str(page_num), font_size=10, align=PP_ALIGN.RIGHT,
                color=RGBColor(100, 100, 100))

    # ── CONTAINER ──
    container = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, C_LEFT, C_TOP, C_WIDTH, C_HEIGHT)
    container.fill.solid()
    container.fill.fore_color.rgb = WHITE
    container.line.color.rgb = BLACK
    container.line.width = Pt(0.75)

    # Gray header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, C_LEFT, C_TOP, C_WIDTH, HEADER_H)
    header.fill.solid()
    header.fill.fore_color.rgb = GRAY_BG
    header.line.color.rgb = BLACK
    header.line.width = Pt(0.75)
    tf = header.text_frame
    tf.margin_left = Pt(10)
    p = tf.paragraphs[0]
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = page_title
    run.font.size = Pt(11)
    run.font.name = "맑은 고딕"
    run.font.color.rgb = BLACK

    # Divider
    add_line(slide, DIVIDER_X, C_TOP + HEADER_H, DIVIDER_X, C_TOP + C_HEIGHT)

    # ══════════════════════════════
    #  LEFT COLUMN (Backend)
    # ══════════════════════════════
    y = C_TOP + HEADER_H + Inches(0.40)
    cx = L_CENTER_X
    bw = L_BOX_W

    # 1) HTML page (Document shape)
    add_document_box(slide, cx - bw / 2, y, bw, BOX_H, html_file, font_size=9)
    y += BOX_H
    add_arrow_down(slide, cx, y, y + ARROW_GAP)
    y += ARROW_GAP

    # 2) Button (Ellipse) - optional
    if has_button:
        btn_w = Inches(1.6)
        btn_h = Inches(0.28)
        add_ellipse_btn(slide, cx - btn_w / 2, y, btn_w, btn_h, button_text, font_size=8)
        y += btn_h
        add_arrow_down(slide, cx, y, y + ARROW_GAP)
        y += ARROW_GAP

    # 3) DB table names + icon
    db_text = "\n".join(db_tables)
    n_db = len(db_tables)
    db_text_h = Inches(max(0.22, n_db * 0.14))

    db_text_left = C_LEFT + Inches(0.15)
    db_text_top = y
    add_textbox(slide, db_text_left, db_text_top,
                Inches(1.7), db_text_h,
                db_text, font_size=7, bold=True, align=PP_ALIGN.LEFT)

    # DB icon image (from template)
    db_icon_left = db_text_left + Inches(0.25)
    db_icon_top = db_text_top + db_text_h + Inches(0.02)
    add_db_icon(slide, db_icon_left, db_icon_top)

    # Dashed arrow from DB icon to Mapper XML
    xml_h = calc_box_h(mapper_xmls)
    add_line(slide,
             db_icon_left + Inches(0.55), db_icon_top + Inches(0.25),
             cx - bw / 2, y + xml_h / 2,
             dashed=True)

    # 4) Mapper XML
    add_process_box(slide, cx - bw / 2, y, bw, xml_h,
                    "\n".join(mapper_xmls), font_size=8)
    y += xml_h
    add_arrow_down(slide, cx, y, y + ARROW_GAP)
    y += ARROW_GAP

    # 5) Mapper Java
    java_h = calc_box_h(mapper_javas)
    add_process_box(slide, cx - bw / 2, y, bw, java_h,
                    "\n".join(mapper_javas), font_size=8)
    y += java_h
    add_arrow_down(slide, cx, y, y + ARROW_GAP)
    y += ARROW_GAP

    # 6) DAO
    dao_h = calc_box_h(daos)
    add_process_box(slide, cx - bw / 2, y, bw, dao_h,
                    "\n".join(daos), font_size=8)
    y += dao_h
    add_arrow_down(slide, cx, y, y + ARROW_GAP)
    y += ARROW_GAP

    # 7) Service (Transaction on ALL slides)
    add_textbox(slide, cx - bw / 2 - Inches(1.15), y,
                Inches(1.0), BOX_H,
                "Transaction", font_size=9, bold=True)
    add_process_box(slide, cx - bw / 2, y, bw, BOX_H, service, font_size=8)
    y += BOX_H
    add_arrow_down(slide, cx, y, y + ARROW_GAP)
    y += ARROW_GAP

    # 8) Controller
    ctrl_top = y
    add_process_box(slide, cx - bw / 2, y, bw, BOX_H, controller, font_size=8)
    ctrl_mid_y = y + BOX_H / 2
    y += BOX_H
    add_arrow_down(slide, cx, y, y + ARROW_GAP)
    y += ARROW_GAP

    # 9) Connect circle + label (bottom of left column)
    circle_cy = y + CONNECT_R + Inches(0.02)
    add_connect_circle(slide, cx, circle_cy)
    add_textbox(slide, cx - Inches(0.35), circle_cy + CONNECT_R + Inches(0.08),
                Inches(0.7), Inches(0.20), "Connect", font_size=8)

    # ── REST label: exactly at Controller height, on divider line ──
    if rest_label:
        rest_w = Inches(0.8) if len(rest_label) > 4 else Inches(0.5)
        add_textbox(slide, DIVIDER_X - rest_w / 2, ctrl_mid_y - Inches(0.12),
                    rest_w, Inches(0.25),
                    rest_label, font_size=10, bold=True)

    # ══════════════════════════════
    #  RIGHT COLUMN (Frontend)
    # ══════════════════════════════
    ry = C_TOP + HEADER_H + Inches(0.50)
    rcx = R_CENTER_X
    rbw = R_BOX_W
    r_gap = Inches(0.22)

    # Connect label + circle (top)
    add_textbox(slide, rcx - Inches(0.35), ry,
                Inches(0.7), Inches(0.20), "Connect", font_size=8)
    ry += Inches(0.24)
    add_connect_circle(slide, rcx, ry + CONNECT_R)
    ry += CONNECT_R * 2 + Inches(0.06)
    add_arrow_down(slide, rcx, ry, ry + r_gap)
    ry += r_gap

    # Decision diamond (optional)
    if has_decision:
        dw = Inches(1.5)
        dh = Inches(0.80)
        diamond_left = rcx - dw / 2
        diamond_top = ry

        add_decision_diamond(slide, diamond_left, diamond_top, dw, dh,
                             decision_text, font_size=7)

        # YES: below-left
        add_textbox(slide, rcx - dw / 2 - Inches(0.08),
                    diamond_top + dh + Inches(0.01),
                    Inches(0.38), Inches(0.20),
                    "YES", font_size=8, bold=True)

        # NO: right of diamond
        add_textbox(slide, rcx + dw / 2 + Inches(0.05),
                    diamond_top + dh / 2 - Inches(0.10),
                    Inches(0.30), Inches(0.20),
                    "NO", font_size=8, bold=True)

        # Line to exception handler
        if exception_handler:
            exc_left = rcx + dw / 2 + Inches(0.38)
            exc_w = Inches(2.2)
            exc_top = diamond_top + dh / 2 - Inches(0.14)
            exc_h = Inches(0.28)
            add_line_arrow(slide,
                           rcx + dw / 2, diamond_top + dh / 2,
                           exc_left, exc_top + exc_h / 2)
            add_process_box(slide, exc_left, exc_top, exc_w, exc_h,
                            exception_handler, font_size=7)

        ry = diamond_top + dh
        add_arrow_down(slide, rcx, ry, ry + r_gap)
        ry += r_gap

    # service.js
    add_process_box(slide, rcx - rbw / 2, ry, rbw, BOX_H, "service.js", font_size=9)
    ry += BOX_H
    add_arrow_down(slide, rcx, ry, ry + r_gap)
    ry += r_gap

    # layout.js
    add_process_box(slide, rcx - rbw / 2, ry, rbw, BOX_H, "layout.js", font_size=9)
    ry += BOX_H
    add_arrow_down(slide, rcx, ry, ry + r_gap)
    ry += r_gap

    # event.js
    add_process_box(slide, rcx - rbw / 2, ry, rbw, BOX_H, "event.js", font_size=9)
    ry += BOX_H
    add_arrow_down(slide, rcx, ry, ry + r_gap)
    ry += r_gap

    # frontend HTML (Document shape)
    add_document_box(slide, rcx - rbw / 2, ry, rbw, BOX_H,
                     frontend_html, font_size=9)

    return slide


slides_data = [
    {
        "page_num": 1,
        "category": "Login Flow",
        "page_title": "Login Page",
        "flow_name": "Login Flow",
        "db_tables": ["tbl_member", "tbl_individual_member", "tbl_oauth"],
        "html_file": "log-in.html",
        "mapper_xmls": ["MemberMapper.xml", "IndividualMemberMapper.xml", "CorpMemberMapper.xml"],
        "mapper_javas": ["MemberMapper.java", "IndividualMemberMapper.java", "CorpMemberMapper.java"],
        "daos": ["MemberDAO.java", "IndividualMemberDAO.java", "CorpMemberDAO.java"],
        "service": "IndividualMemberService.java",
        "controller": "MemberController.java",
        "has_button": True,
        "button_text": "Login Button",
        "frontend_html": "log-in.html",
        "rest_label": "REST",
    },
    {
        "page_num": 2,
        "category": "Login Flow",
        "page_title": "Individual Join Page",
        "flow_name": "Individual\nJoin Flow",
        "db_tables": ["tbl_member", "tbl_individual_member", "tbl_oauth"],
        "html_file": "individual-join.html",
        "mapper_xmls": ["MemberMapper.xml", "IndividualMemberMapper.xml"],
        "mapper_javas": ["MemberMapper.java", "IndividualMemberMapper.java"],
        "daos": ["MemberDAO.java", "IndividualMemberDAO.java"],
        "service": "IndividualMemberService.java",
        "controller": "MemberController.java",
        "has_button": True,
        "button_text": "Join Button",
        "frontend_html": "individual-join.html",
        "has_decision": True,
        "decision_text": "아이디/이메일\n중복 여부",
        "exception_handler": "GlobalRestExceptionHandler.java",
        "rest_label": "REST",
    },
    {
        "page_num": 3,
        "category": "Login Flow",
        "page_title": "Company Join Page",
        "flow_name": "Company\nJoin Flow",
        "db_tables": ["tbl_member", "tbl_corp", "tbl_address"],
        "html_file": "company-join.html",
        "mapper_xmls": ["MemberMapper.xml", "CorpMemberMapper.xml", "AddressMapper.xml"],
        "mapper_javas": ["MemberMapper.java", "CorpMemberMapper.java", "AddressMapper.java"],
        "daos": ["MemberDAO.java", "CorpMemberDAO.java", "AddressDAO.java"],
        "service": "CorpService.java",
        "controller": "MemberController.java",
        "has_button": True,
        "button_text": "Join Button",
        "frontend_html": "company-join.html",
        "has_decision": True,
        "decision_text": "회사명/사업자번호\n중복 여부",
        "exception_handler": "GlobalRestExceptionHandler.java",
        "rest_label": "REST",
    },
    {
        "page_num": 4,
        "category": "MyPage Flow",
        "page_title": "MyPage Main Profile Page",
        "flow_name": "1. Main\nProfile Flow",
        "db_tables": ["view_member_profile", "tbl_latest_watch_posting", "tbl_scrap_posting"],
        "html_file": "mypage.html",
        "mapper_xmls": ["MyPageMapper.xml", "LatestWatchPostingMapper.xml", "ScrapPostingMapper.xml"],
        "mapper_javas": ["MyPageMapper.java", "LatestWatchPostingMapper.java", "ScrapPostingMapper.java"],
        "daos": ["MyPageDAO.java", "LatestWatchPostingDAO.java", "ScrapPostingDAO.java"],
        "service": "MyPageService.java",
        "controller": "MyPageController.java",
        "frontend_html": "mypage.html",
        "rest_label": "REST",
    },
    {
        "page_num": 5,
        "category": "MyPage Flow",
        "page_title": "Profile Edit Page",
        "flow_name": "2. Profile\nEdit Flow",
        "db_tables": ["tbl_member", "tbl_individual_member", "tbl_file"],
        "html_file": "change-my-information.html",
        "mapper_xmls": ["MyPageMapper.xml", "FileMapper.xml"],
        "mapper_javas": ["MyPageMapper.java", "FileMapper.java"],
        "daos": ["MyPageDAO.java", "FileDAO.java"],
        "service": "MyPageService.java",
        "controller": "MyPageController.java",
        "has_button": True,
        "button_text": "Edit Button",
        "frontend_html": "change-my-information.html",
        "rest_label": "REST",
    },
    {
        "page_num": 6,
        "category": "MyPage Flow",
        "page_title": "Account Deletion Page",
        "flow_name": "3. Account\nDeletion Flow",
        "db_tables": ["tbl_member"],
        "html_file": "unsubscribe.html",
        "mapper_xmls": ["MyPageMapper.xml"],
        "mapper_javas": ["MyPageMapper.java"],
        "daos": ["MyPageDAO.java"],
        "service": "MyPageService.java",
        "controller": "MyPageController.java",
        "has_button": True,
        "button_text": "Delete Button",
        "frontend_html": "unsubscribe.html",
        "has_decision": True,
        "decision_text": "입력한 이름\n일치 여부",
        "exception_handler": "UnsubscribeNameMismatch\nException.java",
        "rest_label": "",
    },
    {
        "page_num": 7,
        "category": "MyPage Flow",
        "page_title": "Notification List Page",
        "flow_name": "4. Notification\nList Flow",
        "db_tables": ["tbl_main_notification"],
        "html_file": "notification.html",
        "mapper_xmls": ["MyPageMapper.xml"],
        "mapper_javas": ["MyPageMapper.java"],
        "daos": ["MyPageDAO.java"],
        "service": "MyPageService.java",
        "controller": "MyPageController.java",
        "frontend_html": "notification.html",
        "rest_label": "",
    },
    {
        "page_num": 8,
        "category": "MyPage Flow",
        "page_title": "Activity History Page",
        "flow_name": "5. Activity\nHistory Flow",
        "db_tables": ["tbl_apply", "tbl_experience_program", "tbl_member"],
        "html_file": "experience.html",
        "mapper_xmls": ["MyPageMapper.xml", "ApplyListMapper.xml"],
        "mapper_javas": ["MyPageMapper.java", "ApplyListMapper.java"],
        "daos": ["MyPageDAO.java", "ApplyListDAO.java"],
        "service": "MyPageService.java",
        "controller": "MyPageController.java",
        "has_button": True,
        "button_text": "Search Button",
        "frontend_html": "experience.html",
        "rest_label": "REST",
    },
    {
        "page_num": 9,
        "category": "Header Flow",
        "page_title": "Header Page",
        "flow_name": "Header Flow",
        "db_tables": ["tbl_main_notification"],
        "html_file": "header.html",
        "mapper_xmls": ["IndividualAlramMapper.xml"],
        "mapper_javas": ["IndividualAlramMapper.java"],
        "daos": ["IndividualAlramDAO.java"],
        "service": "IndividualAlramService.java",
        "controller": "IndividualAlramController.java",
        "has_button": True,
        "button_text": "Read Alarm Button",
        "frontend_html": "header.html",
        "rest_label": "REST (PUT)",
    },
    {
        "page_num": 10,
        "category": "Point Flow",
        "page_title": "Point Page",
        "flow_name": "Point Flow",
        "db_tables": ["tbl_point_details", "tbl_individual_member"],
        "html_file": "point.html",
        "mapper_xmls": ["PointDetailsMapper.xml"],
        "mapper_javas": ["PointDetailsMapper.java"],
        "daos": ["PointDetailsDAO.java"],
        "service": "PointService.java",
        "controller": "PointController.java",
        "has_button": True,
        "button_text": "Charge Button",
        "frontend_html": "point.html",
        "has_decision": True,
        "decision_text": "충전/취소\n유효성 검증",
        "exception_handler": "GlobalRestExceptionHandler.java",
        "rest_label": "REST",
    },
]


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for data in slides_data:
        create_slide(
            prs,
            page_num=data["page_num"],
            category=data["category"],
            page_title=data["page_title"],
            flow_name=data["flow_name"],
            db_tables=data["db_tables"],
            html_file=data["html_file"],
            mapper_xmls=data["mapper_xmls"],
            mapper_javas=data["mapper_javas"],
            daos=data["daos"],
            service=data["service"],
            controller=data["controller"],
            has_button=data.get("has_button", False),
            button_text=data.get("button_text", ""),
            frontend_html=data.get("frontend_html"),
            has_decision=data.get("has_decision", False),
            decision_text=data.get("decision_text", ""),
            exception_handler=data.get("exception_handler"),
            rest_label=data.get("rest_label", "REST"),
        )

    out_path = os.path.join(SCRIPT_DIR, "김윤찬_TRY-CATCH_플로우차트.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
